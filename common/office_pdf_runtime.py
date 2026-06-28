"""Shared Office/PDF runtime probes, inspection, render, and QA evidence helpers."""

from __future__ import annotations

import csv
import hashlib
import hmac
import importlib.util
import os
import re
import shutil
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple


SCHEMA_VERSION = 1
PACK_ID = "office-pdf"
_EVIDENCE_HMAC_KEY = b"ecorex-office-pdf-evidence-v1"
_TRUSTED_RENDER_REGISTRY: set[str] = set()
SPREADSHEET_MAX_SHEETS = 12
SPREADSHEET_MAX_SHEET_EVIDENCE = 24
SPREADSHEET_MAX_ROWS_PER_SHEET = 5000
SPREADSHEET_MAX_CELLS = 200000
SPREADSHEET_MAX_COLS_PER_SHEET = 200
CSV_MAX_ROWS = 5000
PRESENTATION_MAX_SLIDES = 120
PRESENTATION_MAX_SLIDE_EVIDENCE = 24
PRESENTATION_OVERLAP_RATIO_THRESHOLD = 0.08
PDF_MAX_PAGES = 120
PDF_MAX_PAGE_EVIDENCE = 24
FORMULA_ERROR_TOKENS = ("#REF!", "#DIV/0!", "#VALUE!", "#NAME?", "#N/A", "#NUM!", "#NULL!")
PDF_GLYPH_RISK_CHARS = ("\ufffd", "\u25a0", "\u25a1", "\x00")


ARTIFACT_KINDS: Dict[str, Dict[str, Any]] = {
    "presentation": {
        "compatibilityId": "office-presentations",
        "officialSkill": "Presentations",
        "extensions": [".ppt", ".pptx"],
        "parserModules": ["pptx"],
        "writerModules": ["pptx"],
        "renderBackends": ["artifact-tool", "libreoffice"],
        "qualityGates": [
            "story-flow",
            "artifact-tool-authoring",
            "layout-bounds",
            "font-size-check",
            "chart-integrity",
            "render-preview",
            "overlap-check",
            "visual-inspection",
        ],
    },
    "spreadsheet": {
        "compatibilityId": "office-spreadsheets",
        "officialSkill": "Spreadsheets",
        "extensions": [".csv", ".tsv", ".xlsm", ".xlsx"],
        "parserModules": ["openpyxl"],
        "writerModules": ["xlsxwriter", "openpyxl"],
        "renderBackends": ["artifact-tool", "libreoffice"],
        "qualityGates": [
            "typed-values",
            "formula-audit",
            "dashboard-structure",
            "chart-render",
            "render-preview",
            "visual-inspection",
            "export-verify",
        ],
    },
    "document": {
        "compatibilityId": "office-documents",
        "officialSkill": "documents",
        "extensions": [".doc", ".docx"],
        "parserModules": ["docx"],
        "writerModules": ["docx"],
        "renderBackends": ["artifact-tool", "libreoffice"],
        "qualityGates": [
            "design-preset",
            "structure-check",
            "render-docx",
            "table-geometry",
            "visual-inspection",
            "redline-preserve",
        ],
    },
    "pdf": {
        "compatibilityId": "office-pdf",
        "officialSkill": "pdf",
        "extensions": [".pdf"],
        "parserModules": ["pypdf", "pdfminer"],
        "writerModules": ["reportlab"],
        "renderModules": ["fitz"],
        "renderBackends": ["pymupdf", "poppler"],
        "qualityGates": [
            "text-orientation",
            "page-render",
            "layout-inspection",
            "table-structure",
            "generation-verify",
            "visual-diff",
        ],
    },
}

_EXTENSION_TO_KIND = {
    extension: kind
    for kind, spec in ARTIFACT_KINDS.items()
    for extension in spec["extensions"]
}


class OfficePdfRuntimeError(RuntimeError):
    """Raised when a requested Office/PDF primitive cannot run in this runtime."""


def detect_artifact_kind(path_or_name: str | os.PathLike[str]) -> str:
    """Return the supported artifact kind for a path or file name, or ``unknown``."""

    extension = Path(str(path_or_name)).suffix.lower()
    return _EXTENSION_TO_KIND.get(extension, "unknown")


def default_quality_gates(kind: str) -> List[str]:
    spec = ARTIFACT_KINDS.get(kind)
    return list(spec.get("qualityGates") or []) if spec else []


def _module_available(name: str) -> bool:
    try:
        return importlib.util.find_spec(name) is not None
    except Exception:
        return False


def _command_available(name: str) -> bool:
    return bool(shutil.which(name))


def _artifact_tool_available() -> bool:
    candidates = [
        Path.cwd() / "node_modules" / "@oai" / "artifact-tool",
        Path(__file__).resolve().parents[1] / "node_modules" / "@oai" / "artifact-tool",
    ]
    env_roots = os.environ.get("NODE_REPL_NODE_MODULE_DIRS") or os.environ.get("ECOREX_NODE_MODULE_DIRS") or ""
    for root in [item for item in env_roots.split(os.pathsep) if item]:
        candidates.append(Path(root) / "@oai" / "artifact-tool")
    return any(path.exists() for path in candidates)


def _hash_ref(value: Any) -> str:
    raw = str(value or "").strip()
    if not raw:
        return ""
    digest = hmac.new(_EVIDENCE_HMAC_KEY, raw.encode("utf-8", errors="replace"), hashlib.sha256).hexdigest()
    return f"hmac:{digest[:16]}"


def _status_from_missing(required: Iterable[str], available: Dict[str, bool]) -> str:
    required_items = list(required)
    if not required_items:
        return "not-applicable"
    missing = [item for item in required_items if not available.get(item)]
    if not missing:
        return "ready"
    if len(missing) == len(required_items):
        return "missing"
    return "partial"


def probe_office_pdf_runtime() -> Dict[str, Any]:
    """Return a redacted readiness snapshot for the shared Office/PDF runtime."""

    module_names = sorted({
        module
        for spec in ARTIFACT_KINDS.values()
        for module in (
            list(spec.get("parserModules") or [])
            + list(spec.get("writerModules") or [])
            + list(spec.get("renderModules") or [])
        )
    } | {"markdownify", "PIL"})
    modules = {name: _module_available(name) for name in module_names}
    commands = {
        "pdfinfo": _command_available("pdfinfo"),
        "pdftoppm": _command_available("pdftoppm"),
        "pdftocairo": _command_available("pdftocairo"),
        "soffice": _command_available("soffice") or _command_available("libreoffice"),
    }
    artifact_tool = _artifact_tool_available()

    kinds: Dict[str, Any] = {}
    for kind, spec in ARTIFACT_KINDS.items():
        parser_modules = list(spec.get("parserModules") or [])
        writer_modules = list(spec.get("writerModules") or [])
        render_modules = list(spec.get("renderModules") or [])
        parse_status = _status_from_missing(parser_modules, modules)
        write_status = _status_from_missing(writer_modules, modules)
        if kind == "pdf":
            render_ready = modules.get("fitz") or commands["pdftoppm"] or commands["pdftocairo"]
            render_status = "ready" if render_ready else "missing"
        else:
            render_ready = artifact_tool or commands["soffice"]
            render_status = "ready" if render_ready else "missing"
        kinds[kind] = {
            "compatibilityId": spec["compatibilityId"],
            "officialSkill": spec["officialSkill"],
            "extensions": list(spec["extensions"]),
            "qualityGates": list(spec["qualityGates"]),
            "parseStatus": parse_status,
            "writeStatus": write_status,
            "renderStatus": render_status,
            "missingParserModules": [item for item in parser_modules if not modules.get(item)],
            "missingWriterModules": [item for item in writer_modules if not modules.get(item)],
            "missingRenderModules": [item for item in render_modules if not modules.get(item)],
            "renderBackends": list(spec["renderBackends"]),
        }

    parse_ready = all(item["parseStatus"] == "ready" for item in kinds.values())
    write_ready = all(item["writeStatus"] in {"ready", "not-applicable"} for item in kinds.values())
    render_ready = all(item["renderStatus"] == "ready" for item in kinds.values())
    return {
        "schemaVersion": SCHEMA_VERSION,
        "packId": PACK_ID,
        "status": "ready" if parse_ready and write_ready and render_ready else "partial",
        "generatedAt": _now(),
        "modules": modules,
        "commands": commands,
        "artifactToolAvailable": artifact_tool,
        "artifactKinds": kinds,
        "redacted": True,
    }


def inspect_office_pdf_artifact(path: str | os.PathLike[str], kind: Optional[str] = None) -> Dict[str, Any]:
    """Inspect an Office/PDF file and return content-free metadata and counts."""

    source = Path(path)
    detected_kind = kind or detect_artifact_kind(source)
    payload: Dict[str, Any] = {
        "schemaVersion": SCHEMA_VERSION,
        "kind": detected_kind,
        "fileRef": _hash_ref(source),
        "extension": source.suffix.lower(),
        "exists": source.exists(),
        "sizeBytes": source.stat().st_size if source.exists() else 0,
        "summary": {},
        "checks": [],
        "redacted": True,
    }
    _add_check(payload["checks"], "file-exists", source.exists())
    if not source.exists():
        return payload
    if detected_kind == "pdf":
        payload["summary"] = _inspect_pdf(source)
    elif detected_kind == "document":
        payload["summary"] = _inspect_docx(source)
    elif detected_kind == "spreadsheet":
        payload["summary"] = _inspect_spreadsheet(source)
    elif detected_kind == "presentation":
        payload["summary"] = _inspect_presentation(source)
    else:
        _add_check(payload["checks"], "supported-extension", False)
    return payload


def render_pdf_pages(
    pdf_path: str | os.PathLike[str],
    output_dir: str | os.PathLike[str],
    *,
    max_pages: int = 4,
    dpi: int = 144,
) -> Dict[str, Any]:
    """Render PDF pages to PNG using PyMuPDF and return redacted render metadata."""

    if not _module_available("fitz"):
        raise OfficePdfRuntimeError("PyMuPDF render backend is not available")
    import fitz  # type: ignore

    source = Path(pdf_path)
    target_dir = Path(output_dir)
    target_dir.mkdir(parents=True, exist_ok=True)
    rendered: List[Dict[str, Any]] = []
    source_ref = _hash_ref(source)
    with fitz.open(str(source)) as doc:
        page_count = len(doc)
        limit = max(0, min(int(max_pages), page_count))
        zoom = max(0.5, float(dpi) / 72.0)
        matrix = fitz.Matrix(zoom, zoom)
        for page_index in range(limit):
            page = doc.load_page(page_index)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            output = target_dir / f"page-{page_index + 1:03d}.png"
            pix.save(str(output))
            item = {
                "page": page_index + 1,
                "artifactRef": _hash_ref(output),
                "extension": ".png",
                "sizeBytes": output.stat().st_size,
                "width": pix.width,
                "height": pix.height,
            }
            item["renderProof"] = _register_trusted_render_artifact(source_ref, item)
            rendered.append(item)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "pdf",
        "sourceRef": source_ref,
        "pageCount": page_count,
        "renderedCount": len(rendered),
        "dpi": dpi,
        "artifacts": rendered,
        "redacted": True,
    }


def render_presentation_preview(
    pptx_path: str | os.PathLike[str],
    output_dir: str | os.PathLike[str],
    *,
    max_slides: int = 4,
    dpi: int = 144,
    timeout_seconds: int = 90,
) -> Dict[str, Any]:
    """Render PPTX slides through LibreOffice-to-PDF plus PDF page rendering."""

    office_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if not office_cmd:
        raise OfficePdfRuntimeError("LibreOffice render backend is not available")
    source = Path(pptx_path)
    target_dir = Path(output_dir)
    pdf_dir = target_dir / "presentation-pdf"
    page_dir = target_dir / "presentation-pages"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    page_dir.mkdir(parents=True, exist_ok=True)
    command = [
        office_cmd,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir),
        str(source),
    ]
    result = subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=max(5, int(timeout_seconds or 90)),
        check=False,
    )
    converted = pdf_dir / f"{source.stem}.pdf"
    if result.returncode != 0 or not converted.exists():
        raise OfficePdfRuntimeError("LibreOffice failed to render presentation preview")
    rendered = render_pdf_pages(converted, page_dir, max_pages=max_slides, dpi=dpi)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "presentation",
        "sourceRef": _hash_ref(source),
        "renderBackend": "libreoffice-pdf",
        "renderedCount": rendered.get("renderedCount", 0),
        "artifacts": rendered.get("artifacts", []),
        "redacted": True,
    }


def render_spreadsheet_preview(
    workbook_path: str | os.PathLike[str],
    output_dir: str | os.PathLike[str],
    *,
    max_sheets: int = 4,
    dpi: int = 144,
    timeout_seconds: int = 90,
) -> Dict[str, Any]:
    """Render spreadsheet sheets through LibreOffice-to-PDF plus PDF page rendering."""

    office_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if not office_cmd:
        raise OfficePdfRuntimeError("LibreOffice render backend is not available")
    source = Path(workbook_path)
    target_dir = Path(output_dir)
    pdf_dir = target_dir / "spreadsheet-pdf"
    page_dir = target_dir / "spreadsheet-pages"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    page_dir.mkdir(parents=True, exist_ok=True)
    command = [
        office_cmd,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir),
        str(source),
    ]
    result = subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=max(5, int(timeout_seconds or 90)),
        check=False,
    )
    converted = pdf_dir / f"{source.stem}.pdf"
    if result.returncode != 0 or not converted.exists():
        raise OfficePdfRuntimeError("LibreOffice failed to render spreadsheet preview")
    rendered = render_pdf_pages(converted, page_dir, max_pages=max_sheets, dpi=dpi)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "spreadsheet",
        "sourceRef": _hash_ref(source),
        "renderBackend": "libreoffice-pdf",
        "renderedCount": rendered.get("renderedCount", 0),
        "artifacts": rendered.get("artifacts", []),
        "redacted": True,
    }


def render_document_preview(
    docx_path: str | os.PathLike[str],
    output_dir: str | os.PathLike[str],
    *,
    max_pages: int = 4,
    dpi: int = 144,
    timeout_seconds: int = 90,
) -> Dict[str, Any]:
    """Render DOCX pages through LibreOffice-to-PDF plus PDF page rendering."""

    office_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if not office_cmd:
        raise OfficePdfRuntimeError("LibreOffice render backend is not available")
    source = Path(docx_path)
    target_dir = Path(output_dir)
    pdf_dir = target_dir / "document-pdf"
    page_dir = target_dir / "document-pages"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    page_dir.mkdir(parents=True, exist_ok=True)
    command = [
        office_cmd,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir),
        str(source),
    ]
    result = subprocess.run(
        command,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=max(5, int(timeout_seconds or 90)),
        check=False,
    )
    converted = pdf_dir / f"{source.stem}.pdf"
    if result.returncode != 0 or not converted.exists():
        raise OfficePdfRuntimeError("LibreOffice failed to render document preview")
    rendered = render_pdf_pages(converted, page_dir, max_pages=max_pages, dpi=dpi)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "document",
        "sourceRef": _hash_ref(source),
        "renderBackend": "libreoffice-pdf",
        "renderedCount": rendered.get("renderedCount", 0),
        "artifacts": rendered.get("artifacts", []),
        "redacted": True,
    }


def build_quality_evidence(
    kind: str,
    checks: Iterable[Dict[str, Any]],
    *,
    source: str | os.PathLike[str] | None = None,
    renders: Optional[Iterable[Dict[str, Any]]] = None,
) -> Dict[str, Any]:
    """Normalize QA checks into a shared evidence shape for Web/API projection."""

    gate_names = default_quality_gates(kind)
    normalized_checks = [_normalize_check(item, gate_names) for item in checks]
    render_items = _sanitize_render_artifacts(renders or [])
    by_id = {item["id"]: item for item in normalized_checks}
    missing_gates = [gate for gate in gate_names if gate not in by_id]
    statuses = {item["status"] for item in normalized_checks}
    status = "pass"
    if "fail" in statuses or missing_gates:
        status = "fail"
    elif "warn" in statuses:
        status = "warn"
    elif "pending" in statuses:
        status = "pending"
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": kind,
        "sourceRef": _hash_ref(source) if source else "",
        "qualityGates": gate_names,
        "checks": normalized_checks,
        "missingQualityGates": missing_gates,
        "status": status,
        "renderedArtifacts": render_items,
        "redacted": True,
    }


def _sanitize_render_artifacts(
    renders: Iterable[Dict[str, Any]],
    *,
    require_render_proof: bool = False,
    require_registered_proof: bool = False,
    trusted_source_ref: str = "",
) -> List[Dict[str, Any]]:
    sanitized: List[Dict[str, Any]] = []
    for raw in renders:
        if not isinstance(raw, dict):
            continue
        item: Dict[str, Any] = {}
        for key in ("slide", "page", "sizeBytes", "width", "height"):
            value = raw.get(key)
            if isinstance(value, bool):
                continue
            if isinstance(value, (int, float)):
                item[key] = max(0, int(value))
        extension = str(raw.get("extension") or "").strip().lower()
        if extension in {".png", ".jpg", ".jpeg", ".webp", ".pdf"}:
            item["extension"] = extension
        for key in ("artifactRef", "sourceRef"):
            value = str(raw.get(key) or "").strip()
            if not value:
                continue
            item[key] = value if _is_hmac_ref(value) else _hash_ref(value)
        proof = str(raw.get("renderProof") or "").strip()
        proof_valid = proof and _is_valid_render_artifact_proof(item, proof)
        proof_registered = (
            proof_valid
            and (
                not require_registered_proof
                or _is_registered_trusted_render_artifact(trusted_source_ref, item, proof)
            )
        )
        if proof_registered:
            item["renderProof"] = proof
        has_ref = bool(item.get("artifactRef") or item.get("sourceRef"))
        has_required_proof = not require_render_proof or bool(item.get("renderProof"))
        if has_ref and has_required_proof:
            sanitized.append(item)
    return sanitized


def _render_artifact_proof(item: Dict[str, Any]) -> str:
    proof_fields = {
        key: item.get(key)
        for key in ("slide", "page", "artifactRef", "sourceRef", "extension", "sizeBytes", "width", "height")
        if item.get(key) is not None
    }
    canonical = "|".join(f"{key}={proof_fields[key]}" for key in sorted(proof_fields))
    return _hash_ref(f"render-proof:{canonical}")


def _register_trusted_render_artifact(source_ref: str, item: Dict[str, Any]) -> str:
    proof = _render_artifact_proof(item)
    _TRUSTED_RENDER_REGISTRY.add(_trusted_render_registry_key(source_ref, item, proof))
    return proof


def _is_registered_trusted_render_artifact(source_ref: str, item: Dict[str, Any], proof: str) -> bool:
    if not source_ref:
        return False
    return _trusted_render_registry_key(source_ref, item, proof) in _TRUSTED_RENDER_REGISTRY


def _trusted_render_registry_key(source_ref: str, item: Dict[str, Any], proof: str) -> str:
    registry_fields = {
        key: item.get(key)
        for key in ("slide", "page", "artifactRef", "sourceRef", "extension", "sizeBytes", "width", "height")
        if item.get(key) is not None
    }
    canonical = "|".join(f"{key}={registry_fields[key]}" for key in sorted(registry_fields))
    return _hash_ref(f"trusted-render:{source_ref}:{canonical}:{proof}")


def _is_valid_render_artifact_proof(item: Dict[str, Any], proof: str) -> bool:
    return _is_hmac_ref(proof) and hmac.compare_digest(proof, _render_artifact_proof(item))


def _is_hmac_ref(value: str) -> bool:
    prefix = "hmac:"
    if not value.startswith(prefix):
        return False
    digest = value[len(prefix):]
    return 8 <= len(digest) <= 64 and all(char in "0123456789abcdef" for char in digest.lower())


def analyze_pdf_quality(path: str | os.PathLike[str]) -> Dict[str, Any]:
    """Return content-free PDF page metrics for extraction, render, and diff QA."""

    if not _module_available("pypdf"):
        raise OfficePdfRuntimeError("pypdf is not available")
    from pypdf import PdfReader  # type: ignore

    source = Path(path)
    reader = PdfReader(str(source))
    metadata = reader.metadata or {}
    encrypted = bool(reader.is_encrypted)
    page_count = _safe_len(reader.pages)
    summary: Dict[str, Any] = {
        "pageCount": page_count,
        "pagesInspected": 0,
        "encrypted": encrypted,
        "metadataKeyCount": len(list(metadata.keys())),
        "textExtractablePageCount": 0,
        "emptyTextPageCount": 0,
        "textExtractionErrorPageCount": 0,
        "totalExtractedTextChars": 0,
        "replacementGlyphCount": 0,
        "nullGlyphCount": 0,
        "glyphRiskPageCount": 0,
        "rotatedPageCount": 0,
        "unexpectedRotationCount": 0,
        "portraitPageCount": 0,
        "landscapePageCount": 0,
        "squarePageCount": 0,
        "pageSizeVariantCount": 0,
        "imageObjectCount": 0,
        "drawingObjectCount": 0,
        "horizontalLineCount": 0,
        "verticalLineCount": 0,
        "tableCandidatePageCount": 0,
        "tableTextCandidatePageCount": 0,
        "blankPageRiskCount": 0,
        "imageOnlyPageCount": 0,
        "fitzInspectionAvailable": _module_available("fitz"),
        "truncated": page_count > PDF_MAX_PAGES,
    }
    page_evidence: List[Dict[str, Any]] = []
    if encrypted:
        return {
            "schemaVersion": SCHEMA_VERSION,
            "kind": "pdf",
            "sourceRef": _hash_ref(source),
            "summary": summary,
            "pageEvidence": page_evidence,
            "redacted": True,
        }

    fitz_doc = None
    if _module_available("fitz"):
        try:
            import fitz  # type: ignore

            fitz_doc = fitz.open(str(source))
        except Exception:
            fitz_doc = None
            summary["fitzInspectionAvailable"] = False

    page_size_buckets: set[Tuple[str, str]] = set()
    try:
        for page_index, page in enumerate(reader.pages):
            if page_index >= PDF_MAX_PAGES:
                summary["truncated"] = True
                break
            width, height = _pdf_page_dimensions(page)
            width_bucket = _dimension_bucket(width)
            height_bucket = _dimension_bucket(height)
            page_size_buckets.add((width_bucket, height_bucket))
            orientation = _page_orientation(width, height)
            if orientation == "portrait":
                summary["portraitPageCount"] += 1
            elif orientation == "landscape":
                summary["landscapePageCount"] += 1
            elif orientation == "square":
                summary["squarePageCount"] += 1

            rotation = _pdf_page_rotation(page)
            if rotation:
                summary["rotatedPageCount"] += 1
            if rotation not in {0, 90, 180, 270}:
                summary["unexpectedRotationCount"] += 1

            text = ""
            text_error = False
            try:
                text = page.extract_text() or ""
            except Exception:
                text_error = True
                summary["textExtractionErrorPageCount"] += 1
            stripped_text = text.strip()
            text_chars = len(stripped_text)
            if stripped_text:
                summary["textExtractablePageCount"] += 1
            else:
                summary["emptyTextPageCount"] += 1
            summary["totalExtractedTextChars"] += text_chars
            replacement_glyphs = sum(text.count(char) for char in ("\ufffd", "\u25a0", "\u25a1"))
            null_glyphs = text.count("\x00")
            summary["replacementGlyphCount"] += replacement_glyphs
            summary["nullGlyphCount"] += null_glyphs
            glyph_risk = replacement_glyphs + null_glyphs
            if glyph_risk:
                summary["glyphRiskPageCount"] += 1

            text_lines = [line for line in text.splitlines() if line.strip()]
            table_text_candidate = _pdf_text_has_table_candidate(text_lines)
            if table_text_candidate:
                summary["tableTextCandidatePageCount"] += 1

            image_count = _count_pdf_page_images(page)
            stream_drawing_info = _inspect_pdf_page_drawing_stream(page)
            drawing_count = 0
            horizontal_lines = 0
            vertical_lines = 0
            fitz_page_info: Dict[str, Any] = {}
            if fitz_doc is not None and page_index < len(fitz_doc):
                fitz_page_info = _inspect_pdf_page_with_fitz(fitz_doc.load_page(page_index))
                image_count = max(image_count, int(fitz_page_info.get("imageObjectCount") or 0))
                drawing_count = int(fitz_page_info.get("drawingObjectCount") or 0)
                horizontal_lines = int(fitz_page_info.get("horizontalLineCount") or 0)
                vertical_lines = int(fitz_page_info.get("verticalLineCount") or 0)
            drawing_count = max(drawing_count, int(stream_drawing_info.get("drawingObjectCount") or 0))
            horizontal_lines = max(horizontal_lines, int(stream_drawing_info.get("horizontalLineCount") or 0))
            vertical_lines = max(vertical_lines, int(stream_drawing_info.get("verticalLineCount") or 0))

            has_table_candidate = bool(
                table_text_candidate
                or fitz_page_info.get("tableCandidate")
                or (horizontal_lines >= 2 and vertical_lines >= 2)
            )
            if has_table_candidate:
                summary["tableCandidatePageCount"] += 1
            summary["imageObjectCount"] += image_count
            summary["drawingObjectCount"] += drawing_count
            summary["horizontalLineCount"] += horizontal_lines
            summary["verticalLineCount"] += vertical_lines

            blank_risk = not stripped_text and image_count == 0 and drawing_count == 0
            if blank_risk:
                summary["blankPageRiskCount"] += 1
            if not stripped_text and image_count > 0:
                summary["imageOnlyPageCount"] += 1

            summary["pagesInspected"] += 1
            if len(page_evidence) < PDF_MAX_PAGE_EVIDENCE:
                page_evidence.append({
                    "page": page_index + 1,
                    "pageRef": _hash_ref(f"{source}#{page_index + 1}"),
                    "orientation": orientation,
                    "rotation": rotation,
                    "widthBucket": width_bucket,
                    "heightBucket": height_bucket,
                    "textLengthBucket": _length_bucket(text_chars),
                    "lineCountBucket": _count_bucket(len(text_lines)),
                    "hasText": bool(stripped_text),
                    "textExtractionError": text_error,
                    "glyphRiskCount": glyph_risk,
                    "imageObjectCount": image_count,
                    "drawingObjectCount": drawing_count,
                    "tableCandidate": has_table_candidate,
                    "blankRisk": blank_risk,
                })
    finally:
        if fitz_doc is not None:
            fitz_doc.close()

    summary["pageSizeVariantCount"] = len(page_size_buckets)
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "pdf",
        "sourceRef": _hash_ref(source),
        "summary": summary,
        "pageEvidence": page_evidence,
        "redacted": True,
    }


def compare_pdf_page_quality(
    reference_path: str | os.PathLike[str],
    candidate_path: str | os.PathLike[str],
) -> Dict[str, Any]:
    """Compare two PDFs using redacted page-level layout buckets."""

    reference = analyze_pdf_quality(reference_path)
    candidate = analyze_pdf_quality(candidate_path)
    return _compare_pdf_analyses(reference, candidate)


def build_pdf_quality_evidence(
    path: str | os.PathLike[str],
    *,
    renders: Optional[Iterable[Dict[str, Any]]] = None,
    reference_path: str | os.PathLike[str] | None = None,
    visual_inspection_passed: bool = False,
) -> Dict[str, Any]:
    """Build the R24-08 PDF page-level quality evidence contract."""

    source = Path(path)
    analysis = analyze_pdf_quality(path)
    summary = dict(analysis.get("summary") or {})
    render_items = _sanitize_render_artifacts(
        renders or [],
        require_render_proof=True,
        require_registered_proof=True,
        trusted_source_ref=_hash_ref(source),
    )
    rendered_count = len(render_items)
    page_count = int(summary.get("pageCount") or 0)
    expected_render_count = min(max(page_count, 1), 4) if page_count else 0
    glyph_issues = int(summary.get("replacementGlyphCount") or 0) + int(summary.get("nullGlyphCount") or 0)
    diff_analysis: Optional[Dict[str, Any]] = None
    diff_mismatch_count = 0
    if reference_path is not None:
        diff_analysis = compare_pdf_page_quality(reference_path, path)
        diff_mismatch_count = int((diff_analysis.get("summary") or {}).get("mismatchCount") or 0)

    checks = [
        {
            "id": "text-orientation",
            "status": (
                "pass"
                if (
                    not summary.get("encrypted")
                    and summary.get("textExtractionErrorPageCount") == 0
                    and summary.get("textExtractablePageCount", 0) > 0
                    and glyph_issues == 0
                    and summary.get("unexpectedRotationCount") == 0
                )
                else "fail"
            ),
            "detail": (
                f"text_pages={summary.get('textExtractablePageCount', 0)}; "
                f"empty_text_pages={summary.get('emptyTextPageCount', 0)}; "
                f"extraction_errors={summary.get('textExtractionErrorPageCount', 0)}; "
                f"glyph_issues={glyph_issues}; "
                f"rotation_issues={summary.get('unexpectedRotationCount', 0)}"
            ),
        },
        {
            "id": "page-render",
            "status": "pass" if rendered_count >= expected_render_count else "fail",
            "detail": f"rendered={rendered_count}; expected_min={expected_render_count}",
        },
        {
            "id": "layout-inspection",
            "status": (
                "pass"
                if summary.get("blankPageRiskCount") == 0 and summary.get("unexpectedRotationCount") == 0
                else "fail"
            ),
            "detail": (
                f"blank_pages={summary.get('blankPageRiskCount', 0)}; "
                f"image_only_pages={summary.get('imageOnlyPageCount', 0)}; "
                f"page_size_variants={summary.get('pageSizeVariantCount', 0)}"
            ),
        },
        {
            "id": "table-structure",
            "status": (
                "pass"
                if summary.get("tableCandidatePageCount", 0) == 0
                or summary.get("tableTextCandidatePageCount", 0) >= summary.get("tableCandidatePageCount", 0)
                or summary.get("textExtractablePageCount", 0) > 0
                else "fail"
            ),
            "detail": (
                f"table_candidates={summary.get('tableCandidatePageCount', 0)}; "
                f"table_text_candidates={summary.get('tableTextCandidatePageCount', 0)}"
            ),
        },
        {
            "id": "generation-verify",
            "status": (
                "pass"
                if source.exists() and source.suffix.lower() == ".pdf" and page_count > 0 and not summary.get("encrypted")
                else "fail"
            ),
            "detail": f"generated=verified; page_count={page_count}" if source.exists() else "generated=missing",
        },
        {
            "id": "visual-diff",
            "status": "pass" if diff_mismatch_count == 0 else "fail",
            "detail": (
                "diff=not_applicable"
                if reference_path is None
                else (
                    f"diff_mismatches={diff_mismatch_count}; "
                    f"pages_compared={(diff_analysis or {}).get('summary', {}).get('pagesCompared', 0)}"
                )
            ),
        },
    ]
    evidence = build_quality_evidence("pdf", checks, source=path, renders=render_items)
    evidence["pdfAnalysis"] = analysis
    if diff_analysis is not None:
        evidence["pdfDiffAnalysis"] = diff_analysis
    return evidence


def analyze_spreadsheet_quality(path: str | os.PathLike[str]) -> Dict[str, Any]:
    """Return content-free workbook quality metrics for formula/chart/render QA."""

    source = Path(path)
    extension = source.suffix.lower()
    if extension in {".csv", ".tsv"}:
        csv_summary = _inspect_csv(source)
        rows = int(csv_summary.get("rowCount") or 0)
        cols = int(csv_summary.get("columnCount") or 0)
        cells = int(csv_summary.get("cellsScanned") or 0)
        summary = {
            "sheetCount": 1,
            "sheetsInspected": 1,
            "emptySheetCount": 0 if rows else 1,
            "nonEmptySheetCount": 1 if rows else 0,
            "hiddenSheetCount": 0,
            "cellsScanned": cells,
            "formulaCellCount": 0,
            "formulaErrorTokenCount": 0,
            "errorCellCount": 0,
            "numericCellCount": 0,
            "dateCellCount": 0,
            "booleanCellCount": 0,
            "textCellCount": cells,
            "blankCellCount": 0,
            "numericTextRiskCount": 0,
            "dateTextRiskCount": 0,
            "tableCount": 0,
            "chartCount": 0,
            "chartIssueCount": 0,
            "mergedRangeCount": 0,
            "truncated": bool(csv_summary.get("truncated")),
        }
        return {
            "schemaVersion": SCHEMA_VERSION,
            "kind": "spreadsheet",
            "sourceRef": _hash_ref(source),
            "summary": summary,
            "sheetEvidence": [
                {
                    "sheet": 1,
                    "visible": True,
                    "rowCount": rows,
                    "columnCount": cols,
                    "nonEmpty": rows > 0,
                    "formulaCellCount": 0,
                    "formulaErrorTokenCount": 0,
                    "numericTextRiskCount": 0,
                    "dateTextRiskCount": 0,
                    "chartCount": 0,
                    "chartIssueCount": 0,
                    "tableCount": 0,
                    "truncated": bool(csv_summary.get("truncated")),
                }
            ],
            "redacted": True,
        }

    if not _module_available("openpyxl"):
        raise OfficePdfRuntimeError("openpyxl is not available")
    import openpyxl  # type: ignore

    workbook = openpyxl.load_workbook(str(source), read_only=False, data_only=False)
    try:
        totals = {
            "sheetCount": len(workbook.worksheets),
            "sheetsInspected": 0,
            "emptySheetCount": 0,
            "nonEmptySheetCount": 0,
            "hiddenSheetCount": 0,
            "cellsScanned": 0,
            "formulaCellCount": 0,
            "formulaErrorTokenCount": 0,
            "errorCellCount": 0,
            "numericCellCount": 0,
            "dateCellCount": 0,
            "booleanCellCount": 0,
            "textCellCount": 0,
            "blankCellCount": 0,
            "numericTextRiskCount": 0,
            "dateTextRiskCount": 0,
            "tableCount": 0,
            "chartCount": 0,
            "chartIssueCount": 0,
            "mergedRangeCount": 0,
            "truncated": len(workbook.worksheets) > SPREADSHEET_MAX_SHEETS,
        }
        sheet_evidence: List[Dict[str, Any]] = []
        remaining_cells = SPREADSHEET_MAX_CELLS

        for sheet_index, worksheet in enumerate(workbook.worksheets):
            if sheet_index >= SPREADSHEET_MAX_SHEETS or remaining_cells <= 0:
                totals["truncated"] = True
                break
            totals["sheetsInspected"] += 1
            if worksheet.sheet_state != "visible":
                totals["hiddenSheetCount"] += 1

            max_row = min(int(worksheet.max_row or 0), SPREADSHEET_MAX_ROWS_PER_SHEET)
            max_col = min(int(worksheet.max_column or 0), SPREADSHEET_MAX_COLS_PER_SHEET)
            sheet_counts = {
                "cellsScanned": 0,
                "nonEmptyCellCount": 0,
                "formulaCellCount": 0,
                "formulaErrorTokenCount": 0,
                "errorCellCount": 0,
                "numericTextRiskCount": 0,
                "dateTextRiskCount": 0,
            }

            for row in worksheet.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col):
                if remaining_cells <= 0:
                    totals["truncated"] = True
                    break
                for cell in row:
                    if remaining_cells <= 0:
                        totals["truncated"] = True
                        break
                    remaining_cells -= 1
                    sheet_counts["cellsScanned"] += 1
                    totals["cellsScanned"] += 1
                    value = cell.value
                    if value is None:
                        totals["blankCellCount"] += 1
                        continue
                    sheet_counts["nonEmptyCellCount"] += 1
                    if cell.data_type == "f" or (isinstance(value, str) and value.startswith("=")):
                        formula = str(value)
                        sheet_counts["formulaCellCount"] += 1
                        totals["formulaCellCount"] += 1
                        if any(token in formula.upper() for token in FORMULA_ERROR_TOKENS):
                            sheet_counts["formulaErrorTokenCount"] += 1
                            totals["formulaErrorTokenCount"] += 1
                        continue
                    if cell.data_type == "e":
                        sheet_counts["errorCellCount"] += 1
                        totals["errorCellCount"] += 1
                    elif cell.is_date:
                        totals["dateCellCount"] += 1
                    elif isinstance(value, bool):
                        totals["booleanCellCount"] += 1
                    elif isinstance(value, (int, float)):
                        totals["numericCellCount"] += 1
                    elif isinstance(value, str):
                        totals["textCellCount"] += 1
                        if _looks_like_numeric_text(value):
                            sheet_counts["numericTextRiskCount"] += 1
                            totals["numericTextRiskCount"] += 1
                        if _looks_like_date_text(value):
                            sheet_counts["dateTextRiskCount"] += 1
                            totals["dateTextRiskCount"] += 1

            charts = list(getattr(worksheet, "_charts", []) or [])
            chart_count = len(charts)
            chart_issue_count = sum(1 for chart in charts if not _spreadsheet_chart_has_data(chart))
            table_count = len(getattr(worksheet, "tables", {}) or {})
            merged_range_count = len(list(getattr(worksheet.merged_cells, "ranges", []) or []))
            totals["chartCount"] += chart_count
            totals["chartIssueCount"] += chart_issue_count
            totals["tableCount"] += table_count
            totals["mergedRangeCount"] += merged_range_count

            non_empty = sheet_counts["nonEmptyCellCount"] > 0 or chart_count > 0 or table_count > 0
            if non_empty:
                totals["nonEmptySheetCount"] += 1
            else:
                totals["emptySheetCount"] += 1

            if len(sheet_evidence) < SPREADSHEET_MAX_SHEET_EVIDENCE:
                sheet_evidence.append({
                    "sheet": sheet_index + 1,
                    "sheetRef": _hash_ref(worksheet.title),
                    "visible": worksheet.sheet_state == "visible",
                    "rowCount": int(worksheet.max_row or 0),
                    "columnCount": int(worksheet.max_column or 0),
                    "nonEmpty": non_empty,
                    "formulaCellCount": sheet_counts["formulaCellCount"],
                    "formulaErrorTokenCount": sheet_counts["formulaErrorTokenCount"],
                    "errorCellCount": sheet_counts["errorCellCount"],
                    "numericTextRiskCount": sheet_counts["numericTextRiskCount"],
                    "dateTextRiskCount": sheet_counts["dateTextRiskCount"],
                    "chartCount": chart_count,
                    "chartIssueCount": chart_issue_count,
                    "tableCount": table_count,
                    "mergedRangeCount": merged_range_count,
                    "truncated": (
                        int(worksheet.max_row or 0) > SPREADSHEET_MAX_ROWS_PER_SHEET
                        or int(worksheet.max_column or 0) > SPREADSHEET_MAX_COLS_PER_SHEET
                    ),
                })

        return {
            "schemaVersion": SCHEMA_VERSION,
            "kind": "spreadsheet",
            "sourceRef": _hash_ref(source),
            "summary": totals,
            "sheetEvidence": sheet_evidence,
            "redacted": True,
        }
    finally:
        workbook.close()


def build_spreadsheet_quality_evidence(
    path: str | os.PathLike[str],
    *,
    renders: Optional[Iterable[Dict[str, Any]]] = None,
    visual_inspection_passed: bool = False,
) -> Dict[str, Any]:
    """Build the R24-06 spreadsheet quality evidence contract."""

    source = Path(path)
    analysis = analyze_spreadsheet_quality(path)
    summary = dict(analysis.get("summary") or {})
    render_items = _sanitize_render_artifacts(renders or [], require_render_proof=True)
    sheet_count = int(summary.get("sheetCount") or 0)
    rendered_count = len(render_items)
    expected_render_count = min(max(sheet_count, 1), 4)
    has_charts = int(summary.get("chartCount") or 0) > 0

    checks = [
        {
            "id": "typed-values",
            "status": (
                "pass"
                if summary.get("numericTextRiskCount") == 0 and summary.get("dateTextRiskCount") == 0
                else "fail"
            ),
            "detail": (
                f"numeric_text={summary.get('numericTextRiskCount', 0)}; "
                f"date_text={summary.get('dateTextRiskCount', 0)}"
            ),
        },
        {
            "id": "formula-audit",
            "status": (
                "pass"
                if summary.get("formulaErrorTokenCount") == 0 and summary.get("errorCellCount") == 0
                else "fail"
            ),
            "detail": (
                f"formulas={summary.get('formulaCellCount', 0)}; "
                f"formula_errors={summary.get('formulaErrorTokenCount', 0)}; "
                f"error_cells={summary.get('errorCellCount', 0)}"
            ),
        },
        {
            "id": "dashboard-structure",
            "status": (
                "pass"
                if sheet_count > 0 and summary.get("nonEmptySheetCount", 0) > 0 and summary.get("emptySheetCount") == 0
                else "fail"
            ),
            "detail": (
                f"sheets={sheet_count}; non_empty_sheets={summary.get('nonEmptySheetCount', 0)}; "
                f"empty_sheets={summary.get('emptySheetCount', 0)}"
            ),
        },
        {
            "id": "chart-render",
            "status": (
                "pass"
                if summary.get("chartIssueCount") == 0 and (not has_charts or rendered_count > 0)
                else "fail"
            ),
            "detail": f"charts={summary.get('chartCount', 0)}; chart_issues={summary.get('chartIssueCount', 0)}",
        },
        {
            "id": "render-preview",
            "status": "pass" if rendered_count >= expected_render_count else "fail",
            "detail": f"rendered={rendered_count}; expected_min={expected_render_count}",
        },
        {
            "id": "visual-inspection",
            "status": "pass" if visual_inspection_passed else "pending",
            "detail": "manual_visual_review=pass" if visual_inspection_passed else "manual_visual_review=pending",
        },
        {
            "id": "export-verify",
            "status": "pass" if source.exists() and source.suffix.lower() in {".xlsx", ".xlsm", ".csv", ".tsv"} else "fail",
            "detail": "export=verified" if source.exists() else "export=missing",
        },
    ]
    evidence = build_quality_evidence("spreadsheet", checks, source=path, renders=render_items)
    evidence["spreadsheetAnalysis"] = analysis
    return evidence


def analyze_document_quality(path: str | os.PathLike[str]) -> Dict[str, Any]:
    """Return content-free DOCX quality metrics for structure/table/redline QA."""

    if not _module_available("docx"):
        raise OfficePdfRuntimeError("python-docx is not available")
    import docx  # type: ignore

    source = Path(path)
    document = docx.Document(str(source))
    ooxml = _inspect_docx_ooxml(source)
    heading_count = 0
    title_count = 0
    list_count = 0
    empty_paragraph_count = 0
    for paragraph in document.paragraphs:
        text_length = len(str(paragraph.text or "").strip())
        if text_length == 0:
            empty_paragraph_count += 1
        style_name = str(getattr(paragraph.style, "name", "") or "").lower()
        if style_name.startswith("heading"):
            heading_count += 1
        if style_name == "title":
            title_count += 1
        if "list" in style_name:
            list_count += 1

    table_cell_count = 0
    table_row_count = 0
    for table in document.tables:
        table_row_count += len(table.rows)
        table_cell_count += sum(len(row.cells) for row in table.rows)

    summary = {
        "paragraphCount": len(document.paragraphs),
        "emptyParagraphCount": empty_paragraph_count,
        "headingParagraphCount": heading_count,
        "titleParagraphCount": title_count,
        "listParagraphCount": list_count,
        "tableCount": len(document.tables),
        "tableRowCount": table_row_count,
        "tableCellCount": table_cell_count,
        "sectionCount": len(document.sections),
        **ooxml,
    }
    summary["trackedChangeCount"] = summary.get("trackedInsertCount", 0) + summary.get("trackedDeleteCount", 0)
    summary["tableIssueCount"] = (
        summary.get("tableGridMissingCount", 0)
        + summary.get("tableCellWidthMissingCount", 0)
        + summary.get("tableWidthMissingCount", 0)
    )
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "document",
        "sourceRef": _hash_ref(source),
        "summary": summary,
        "documentEvidence": {
            "hasTitleOrHeading": title_count > 0 or heading_count > 0,
            "hasBody": len(document.paragraphs) > empty_paragraph_count or len(document.tables) > 0,
            "hasTables": len(document.tables) > 0,
            "hasComments": summary.get("commentCount", 0) > 0,
            "hasTrackedChanges": summary.get("trackedChangeCount", 0) > 0,
        },
        "redacted": True,
    }


def build_document_quality_evidence(
    path: str | os.PathLike[str],
    *,
    renders: Optional[Iterable[Dict[str, Any]]] = None,
    visual_inspection_passed: bool = False,
) -> Dict[str, Any]:
    """Build the R24-07 document quality evidence contract."""

    analysis = analyze_document_quality(path)
    summary = dict(analysis.get("summary") or {})
    render_items = _sanitize_render_artifacts(renders or [], require_render_proof=True)
    rendered_count = len(render_items)
    expected_render_count = 1 if summary.get("paragraphCount", 0) or summary.get("tableCount", 0) else 0
    comment_mismatch = (
        int(summary.get("commentCount") or 0) != int(summary.get("commentReferenceCount") or 0)
        or int(summary.get("commentIdMismatchCount") or 0) > 0
    )

    checks = [
        {
            "id": "design-preset",
            "status": "pass" if summary.get("titleParagraphCount", 0) + summary.get("headingParagraphCount", 0) > 0 else "fail",
            "detail": (
                f"titles={summary.get('titleParagraphCount', 0)}; "
                f"headings={summary.get('headingParagraphCount', 0)}"
            ),
        },
        {
            "id": "structure-check",
            "status": "pass" if summary.get("sectionCount", 0) > 0 and summary.get("paragraphCount", 0) > 0 else "fail",
            "detail": (
                f"paragraphs={summary.get('paragraphCount', 0)}; "
                f"sections={summary.get('sectionCount', 0)}"
            ),
        },
        {
            "id": "render-docx",
            "status": "pass" if rendered_count >= expected_render_count else "fail",
            "detail": f"rendered={rendered_count}; expected_min={expected_render_count}",
        },
        {
            "id": "table-geometry",
            "status": "pass" if summary.get("tableIssueCount", 0) == 0 else "fail",
            "detail": (
                f"tables={summary.get('tableCount', 0)}; "
                f"table_issues={summary.get('tableIssueCount', 0)}"
            ),
        },
        {
            "id": "visual-inspection",
            "status": "pass" if visual_inspection_passed else "pending",
            "detail": "manual_visual_review=pass" if visual_inspection_passed else "manual_visual_review=pending",
        },
        {
            "id": "redline-preserve",
            "status": "pass" if not comment_mismatch else "fail",
            "detail": (
                f"comments={summary.get('commentCount', 0)}; "
                f"comment_refs={summary.get('commentReferenceCount', 0)}; "
                f"comment_id_mismatches={summary.get('commentIdMismatchCount', 0)}; "
                f"tracked_changes={summary.get('trackedChangeCount', 0)}"
            ),
        },
    ]
    evidence = build_quality_evidence("document", checks, source=path, renders=render_items)
    evidence["documentAnalysis"] = analysis
    return evidence


def analyze_presentation_quality(path: str | os.PathLike[str]) -> Dict[str, Any]:
    """Return content-free PPTX quality metrics for story/layout/chart QA."""

    if not _module_available("pptx"):
        raise OfficePdfRuntimeError("python-pptx is not available")
    from pptx import Presentation  # type: ignore

    source = Path(path)
    presentation = Presentation(str(source))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    slide_area = max(1, slide_width * slide_height)

    slide_evidence: List[Dict[str, Any]] = []
    totals = {
        "slideCount": len(presentation.slides),
        "slidesInspected": 0,
        "emptySlideCount": 0,
        "missingTitleCount": 0,
        "titleWrapRiskCount": 0,
        "outOfBoundsCount": 0,
        "overlapWarningCount": 0,
        "fontViolationCount": 0,
        "fontUnspecifiedRunCount": 0,
        "chartCount": 0,
        "chartIssueCount": 0,
        "pictureCount": 0,
        "tableCount": 0,
        "textShapeCount": 0,
        "truncated": len(presentation.slides) > PRESENTATION_MAX_SLIDES,
    }

    for slide_index, slide in enumerate(presentation.slides):
        if slide_index >= PRESENTATION_MAX_SLIDES:
            break
        slide_no = slide_index + 1
        totals["slidesInspected"] += 1
        title_shape = slide.shapes.title
        title_text_len = _shape_text_length(title_shape) if title_shape is not None else 0
        if title_text_len <= 0:
            totals["missingTitleCount"] += 1
        if title_text_len > 70:
            totals["titleWrapRiskCount"] += 1

        bounds: List[Tuple[int, int, int, int, int, str]] = []
        visible_payload = False
        slide_out_of_bounds = 0
        slide_font_violations = 0
        slide_unspecified_fonts = 0
        slide_chart_issues = 0
        slide_chart_count = 0
        slide_text_shape_count = 0

        for shape_index, shape in enumerate(slide.shapes):
            if _shape_has_visible_payload(shape):
                visible_payload = True
            if getattr(shape, "has_chart", False):
                totals["chartCount"] += 1
                slide_chart_count += 1
                if not _chart_has_data(shape):
                    totals["chartIssueCount"] += 1
                    slide_chart_issues += 1
            if getattr(shape, "has_table", False):
                totals["tableCount"] += 1
            if _shape_type_name(shape) == "PICTURE":
                totals["pictureCount"] += 1
            if getattr(shape, "has_text_frame", False):
                totals["textShapeCount"] += 1
                slide_text_shape_count += 1
                role = _presentation_text_role(shape, title_shape, slide_index)
                violations, unspecified = _font_check_counts(shape, role)
                totals["fontViolationCount"] += violations
                totals["fontUnspecifiedRunCount"] += unspecified
                slide_font_violations += violations
                slide_unspecified_fonts += unspecified

            shape_bounds = _shape_bounds(shape)
            if not shape_bounds:
                continue
            left, top, width, height = shape_bounds
            if left < 0 or top < 0 or left + width > slide_width or top + height > slide_height:
                totals["outOfBoundsCount"] += 1
                slide_out_of_bounds += 1
            if _is_background_like_shape(shape_bounds, slide_width, slide_height, slide_area):
                continue
            bounds.append((left, top, width, height, shape_index, _shape_type_name(shape)))

        overlap_count = _count_shape_overlaps(bounds)
        totals["overlapWarningCount"] += overlap_count
        if not visible_payload:
            totals["emptySlideCount"] += 1

        if len(slide_evidence) < PRESENTATION_MAX_SLIDE_EVIDENCE:
            slide_evidence.append({
                "slide": slide_no,
                "shapeCount": len(slide.shapes),
                "textShapeCount": slide_text_shape_count,
                "titlePresent": title_text_len > 0,
                "titleLengthBucket": _length_bucket(title_text_len),
                "titleWrapRisk": title_text_len > 70,
                "empty": not visible_payload,
                "outOfBoundsCount": slide_out_of_bounds,
                "overlapWarningCount": overlap_count,
                "fontViolationCount": slide_font_violations,
                "fontUnspecifiedRunCount": slide_unspecified_fonts,
                "chartCount": slide_chart_count,
                "chartIssueCount": slide_chart_issues,
            })

    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "presentation",
        "sourceRef": _hash_ref(source),
        "slideSize": {"widthEmu": slide_width, "heightEmu": slide_height},
        "summary": totals,
        "slideEvidence": slide_evidence,
        "redacted": True,
    }


def build_presentation_quality_evidence(
    path: str | os.PathLike[str],
    *,
    authoring_route: str = "",
    renders: Optional[Iterable[Dict[str, Any]]] = None,
    visual_inspection_passed: bool = False,
) -> Dict[str, Any]:
    """Build the R24-05 PPT quality evidence contract from static and render QA."""

    analysis = analyze_presentation_quality(path)
    summary = dict(analysis.get("summary") or {})
    render_items = _sanitize_render_artifacts(renders or [], require_render_proof=True)
    slide_count = int(summary.get("slideCount") or 0)
    rendered_count = len(render_items)
    expected_render_count = min(max(slide_count, 1), 4)
    route = _normalize_presentation_authoring_route(authoring_route)

    checks = [
        {
            "id": "story-flow",
            "status": "pass" if slide_count > 0 and summary.get("missingTitleCount") == 0 and summary.get("emptySlideCount") == 0 else "fail",
            "detail": (
                f"slides={slide_count}; missing_titles={summary.get('missingTitleCount', 0)}; "
                f"empty_slides={summary.get('emptySlideCount', 0)}"
            ),
        },
        {
            "id": "artifact-tool-authoring",
            "status": "pass" if route != "unspecified" else "pending",
            "detail": f"route={route}",
        },
        {
            "id": "layout-bounds",
            "status": "pass" if summary.get("outOfBoundsCount") == 0 else "fail",
            "detail": f"out_of_bounds={summary.get('outOfBoundsCount', 0)}",
        },
        {
            "id": "font-size-check",
            "status": "pass" if summary.get("fontViolationCount") == 0 else "fail",
            "detail": (
                f"violations={summary.get('fontViolationCount', 0)}; "
                f"unspecified={summary.get('fontUnspecifiedRunCount', 0)}"
            ),
        },
        {
            "id": "chart-integrity",
            "status": "pass" if summary.get("chartIssueCount") == 0 else "fail",
            "detail": f"charts={summary.get('chartCount', 0)}; issues={summary.get('chartIssueCount', 0)}",
        },
        {
            "id": "render-preview",
            "status": "pass" if rendered_count >= expected_render_count else "fail",
            "detail": f"rendered={rendered_count}; expected_min={expected_render_count}",
        },
        {
            "id": "overlap-check",
            "status": "pass" if summary.get("overlapWarningCount") == 0 else "fail",
            "detail": f"overlaps={summary.get('overlapWarningCount', 0)}",
        },
        {
            "id": "visual-inspection",
            "status": "pass" if visual_inspection_passed else "pending",
            "detail": "manual_visual_review=pass" if visual_inspection_passed else "manual_visual_review=pending",
        },
    ]
    evidence = build_quality_evidence("presentation", checks, source=path, renders=render_items)
    evidence["presentationAnalysis"] = analysis
    evidence["authoringRoute"] = route
    return evidence


def _normalize_presentation_authoring_route(value: str) -> str:
    route = str(value or "").strip().lower()
    return route if route in {"artifact-tool", "template-following", "verified-existing-deck"} else "unspecified"


def _safe_len(value: Any) -> int:
    try:
        return len(value)
    except Exception:
        return 0


def _pdf_page_dimensions(page: Any) -> Tuple[int, int]:
    try:
        box = page.mediabox
        return max(0, int(float(box.width))), max(0, int(float(box.height)))
    except Exception:
        return 0, 0


def _pdf_page_rotation(page: Any) -> int:
    try:
        return int(page.get("/Rotate", 0) or 0) % 360
    except Exception:
        return 0


def _dimension_bucket(value: int) -> str:
    numeric = max(0, int(value or 0))
    if numeric <= 0:
        return "unknown"
    rounded = int(round(numeric / 25.0) * 25)
    return str(max(25, rounded))


def _count_bucket(value: int) -> str:
    numeric = max(0, int(value or 0))
    if numeric == 0:
        return "none"
    if numeric <= 3:
        return "few"
    if numeric <= 12:
        return "some"
    return "many"


def _page_orientation(width: int, height: int) -> str:
    if width <= 0 or height <= 0:
        return "unknown"
    if abs(width - height) <= max(width, height) * 0.03:
        return "square"
    return "landscape" if width > height else "portrait"


def _pdf_text_has_table_candidate(lines: Iterable[str]) -> bool:
    candidate_lines = 0
    for line in lines:
        text = str(line or "").strip()
        if not text:
            continue
        if re.search(r"\S\s{2,}\S", text) or len(text.split()) >= 4:
            candidate_lines += 1
        if candidate_lines >= 2:
            return True
    return False


def _count_pdf_page_images(page: Any) -> int:
    try:
        resources = page.get("/Resources") or {}
        if hasattr(resources, "get_object"):
            resources = resources.get_object()
        xobjects = resources.get("/XObject") if hasattr(resources, "get") else None
        if hasattr(xobjects, "get_object"):
            xobjects = xobjects.get_object()
        if not hasattr(xobjects, "values"):
            return 0
        count = 0
        for raw_object in xobjects.values():
            try:
                obj = raw_object.get_object() if hasattr(raw_object, "get_object") else raw_object
                if obj.get("/Subtype") == "/Image":
                    count += 1
            except Exception:
                continue
        return count
    except Exception:
        return 0


def _inspect_pdf_page_drawing_stream(page: Any) -> Dict[str, int]:
    try:
        contents = page.get_contents()
        if contents is None:
            data = b""
        elif isinstance(contents, list):
            data = b"".join(item.get_data() for item in contents if hasattr(item, "get_data"))
        elif hasattr(contents, "get_data"):
            data = contents.get_data()
        else:
            data = b""
    except Exception:
        data = b""
    if not data:
        return {"drawingObjectCount": 0, "horizontalLineCount": 0, "verticalLineCount": 0}
    text = data.decode("latin-1", errors="ignore")
    rectangle_count = len(re.findall(r"(?:^|\s)re(?:\s|$)", text))
    line_count = len(re.findall(r"(?:^|\s)l(?:\s|$)", text))
    return {
        "drawingObjectCount": rectangle_count + line_count,
        "horizontalLineCount": rectangle_count * 2,
        "verticalLineCount": rectangle_count * 2,
    }


def _inspect_pdf_page_with_fitz(page: Any) -> Dict[str, Any]:
    drawing_count = 0
    horizontal_lines = 0
    vertical_lines = 0
    image_count = 0
    try:
        blocks = page.get_text("dict").get("blocks", [])
        image_count = sum(1 for block in blocks if block.get("type") == 1)
    except Exception:
        image_count = 0
    try:
        drawings = page.get_drawings()
    except Exception:
        drawings = []
    for drawing in drawings:
        items = drawing.get("items") or []
        drawing_count += len(items)
        for item in items:
            if not isinstance(item, (tuple, list)) or not item:
                continue
            op = item[0]
            if op == "l" and len(item) >= 3:
                p1, p2 = item[1], item[2]
                x1, y1 = float(getattr(p1, "x", 0.0)), float(getattr(p1, "y", 0.0))
                x2, y2 = float(getattr(p2, "x", 0.0)), float(getattr(p2, "y", 0.0))
                if abs(y1 - y2) <= 1.0 and abs(x1 - x2) >= 8.0:
                    horizontal_lines += 1
                elif abs(x1 - x2) <= 1.0 and abs(y1 - y2) >= 8.0:
                    vertical_lines += 1
            elif op == "re" and len(item) >= 2:
                rect = item[1]
                width = abs(float(getattr(rect, "x1", 0.0)) - float(getattr(rect, "x0", 0.0)))
                height = abs(float(getattr(rect, "y1", 0.0)) - float(getattr(rect, "y0", 0.0)))
                if width >= 8.0 and height >= 8.0:
                    horizontal_lines += 2
                    vertical_lines += 2
    return {
        "imageObjectCount": image_count,
        "drawingObjectCount": drawing_count,
        "horizontalLineCount": horizontal_lines,
        "verticalLineCount": vertical_lines,
        "tableCandidate": horizontal_lines >= 2 and vertical_lines >= 2,
    }


def _compare_pdf_analyses(reference: Dict[str, Any], candidate: Dict[str, Any]) -> Dict[str, Any]:
    reference_summary = dict(reference.get("summary") or {})
    candidate_summary = dict(candidate.get("summary") or {})
    reference_pages = list(reference.get("pageEvidence") or [])
    candidate_pages = list(candidate.get("pageEvidence") or [])
    pages_compared = min(len(reference_pages), len(candidate_pages))
    page_size_mismatches = 0
    orientation_mismatches = 0
    text_bucket_mismatches = 0
    table_candidate_mismatches = 0
    image_count_mismatches = 0
    page_diff_evidence: List[Dict[str, Any]] = []

    for index in range(pages_compared):
        ref_page = reference_pages[index]
        candidate_page = candidate_pages[index]
        size_match = (
            ref_page.get("widthBucket") == candidate_page.get("widthBucket")
            and ref_page.get("heightBucket") == candidate_page.get("heightBucket")
        )
        orientation_match = ref_page.get("orientation") == candidate_page.get("orientation")
        text_bucket_match = ref_page.get("textLengthBucket") == candidate_page.get("textLengthBucket")
        table_candidate_match = bool(ref_page.get("tableCandidate")) == bool(candidate_page.get("tableCandidate"))
        image_match = int(ref_page.get("imageObjectCount") or 0) == int(candidate_page.get("imageObjectCount") or 0)
        if not size_match:
            page_size_mismatches += 1
        if not orientation_match:
            orientation_mismatches += 1
        if not text_bucket_match:
            text_bucket_mismatches += 1
        if not table_candidate_match:
            table_candidate_mismatches += 1
        if not image_match:
            image_count_mismatches += 1
        if len(page_diff_evidence) < PDF_MAX_PAGE_EVIDENCE:
            page_diff_evidence.append({
                "page": index + 1,
                "status": (
                    "pass"
                    if size_match and orientation_match and text_bucket_match and table_candidate_match and image_match
                    else "fail"
                ),
                "sizeMatch": size_match,
                "orientationMatch": orientation_match,
                "textLengthBucketMatch": text_bucket_match,
                "tableCandidateMatch": table_candidate_match,
                "imageObjectCountMatch": image_match,
            })

    page_count_match = int(reference_summary.get("pageCount") or 0) == int(candidate_summary.get("pageCount") or 0)
    page_count_mismatch = 0 if page_count_match else 1
    mismatch_count = (
        page_count_mismatch
        + page_size_mismatches
        + orientation_mismatches
        + text_bucket_mismatches
        + table_candidate_mismatches
        + image_count_mismatches
    )
    return {
        "schemaVersion": SCHEMA_VERSION,
        "kind": "pdf-diff",
        "referenceRef": str(reference.get("sourceRef") or ""),
        "candidateRef": str(candidate.get("sourceRef") or ""),
        "summary": {
            "referencePageCount": int(reference_summary.get("pageCount") or 0),
            "candidatePageCount": int(candidate_summary.get("pageCount") or 0),
            "pageCountMatch": page_count_match,
            "pagesCompared": pages_compared,
            "pageSizeMismatchCount": page_size_mismatches,
            "orientationMismatchCount": orientation_mismatches,
            "textLengthBucketMismatchCount": text_bucket_mismatches,
            "tableCandidateMismatchCount": table_candidate_mismatches,
            "imageObjectCountMismatchCount": image_count_mismatches,
            "mismatchCount": mismatch_count,
        },
        "pageDiffEvidence": page_diff_evidence,
        "redacted": True,
    }


def _inspect_pdf(path: Path) -> Dict[str, Any]:
    if not _module_available("pypdf"):
        raise OfficePdfRuntimeError("pypdf is not available")
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    metadata = reader.metadata or {}
    return {
        "pageCount": len(reader.pages),
        "encrypted": bool(reader.is_encrypted),
        "metadataKeyCount": len(list(metadata.keys())),
    }


def _inspect_docx_ooxml(path: Path) -> Dict[str, int]:
    counts = {
        "commentCount": 0,
        "commentReferenceCount": 0,
        "commentIdMismatchCount": 0,
        "trackedInsertCount": 0,
        "trackedDeleteCount": 0,
        "tableGridMissingCount": 0,
        "tableWidthMissingCount": 0,
        "tableCellWidthMissingCount": 0,
    }
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    try:
        with zipfile.ZipFile(path) as archive:
            document_xml = archive.read("word/document.xml")
            root = ET.fromstring(document_xml)
            comment_ref_ids = {
                str(item.get(f"{{{namespace['w']}}}id") or "")
                for item in root.findall(".//w:commentReference", namespace)
            }
            comment_ref_ids.discard("")
            counts["commentReferenceCount"] = len(comment_ref_ids)
            counts["trackedInsertCount"] = len(root.findall(".//w:ins", namespace))
            counts["trackedDeleteCount"] = len(root.findall(".//w:del", namespace))
            for table in root.findall(".//w:tbl", namespace):
                if table.find("w:tblGrid", namespace) is None:
                    counts["tableGridMissingCount"] += 1
                properties = table.find("w:tblPr", namespace)
                if properties is None or properties.find("w:tblW", namespace) is None:
                    counts["tableWidthMissingCount"] += 1
                for cell in table.findall(".//w:tc", namespace):
                    cell_properties = cell.find("w:tcPr", namespace)
                    if cell_properties is None or cell_properties.find("w:tcW", namespace) is None:
                        counts["tableCellWidthMissingCount"] += 1
            if "word/comments.xml" in archive.namelist():
                comments_root = ET.fromstring(archive.read("word/comments.xml"))
                comment_ids = {
                    str(item.get(f"{{{namespace['w']}}}id") or "")
                    for item in comments_root.findall(".//w:comment", namespace)
                }
                comment_ids.discard("")
                counts["commentCount"] = len(comment_ids)
                counts["commentIdMismatchCount"] = len(comment_ids.symmetric_difference(comment_ref_ids))
    except Exception:
        pass
    return counts


def _inspect_docx(path: Path) -> Dict[str, Any]:
    if not _module_available("docx"):
        raise OfficePdfRuntimeError("python-docx is not available")
    import docx  # type: ignore

    document = docx.Document(str(path))
    return {
        "paragraphCount": len(document.paragraphs),
        "tableCount": len(document.tables),
        "sectionCount": len(document.sections),
    }


def _inspect_spreadsheet(path: Path) -> Dict[str, Any]:
    extension = path.suffix.lower()
    if extension in {".csv", ".tsv"}:
        return _inspect_csv(path)
    if not _module_available("openpyxl"):
        raise OfficePdfRuntimeError("openpyxl is not available")
    import openpyxl  # type: ignore

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=False)
    try:
        dimensions = []
        formula_cells = 0
        cells_scanned = 0
        sheets_inspected = 0
        truncated = False
        for sheet_index, worksheet in enumerate(workbook.worksheets):
            if sheet_index >= SPREADSHEET_MAX_SHEETS:
                truncated = True
                break
            sheets_inspected += 1
            dimensions.append({"maxRows": worksheet.max_row, "maxColumns": worksheet.max_column})
            row_limit = min(int(worksheet.max_row or 0), SPREADSHEET_MAX_ROWS_PER_SHEET) or None
            rows_seen = 0
            for row in worksheet.iter_rows(max_row=row_limit):
                rows_seen += 1
                for cell in row:
                    cells_scanned += 1
                    value = cell.value
                    if isinstance(value, str) and value.startswith("="):
                        formula_cells += 1
                    if cells_scanned >= SPREADSHEET_MAX_CELLS:
                        truncated = True
                        break
                if truncated:
                    break
            if int(worksheet.max_row or 0) > rows_seen:
                truncated = True
            if truncated:
                break
        return {
            "sheetCount": len(workbook.worksheets),
            "dimensions": dimensions[:12],
            "formulaCellCount": formula_cells,
            "cellsScanned": cells_scanned,
            "sheetsInspected": sheets_inspected,
            "truncated": truncated,
        }
    finally:
        workbook.close()


def _inspect_csv(path: Path) -> Dict[str, Any]:
    row_count = 0
    max_columns = 0
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.reader(handle, delimiter="\t" if path.suffix.lower() == ".tsv" else ",")
        for row_count, row in enumerate(reader, start=1):
            max_columns = max(max_columns, len(row))
            if row_count >= CSV_MAX_ROWS:
                break
    return {"rowsSampled": row_count, "maxColumns": max_columns, "truncated": row_count >= CSV_MAX_ROWS}


def _inspect_presentation(path: Path) -> Dict[str, Any]:
    if not _module_available("pptx"):
        raise OfficePdfRuntimeError("python-pptx is not available")
    from pptx import Presentation  # type: ignore

    presentation = Presentation(str(path))
    slide_count = len(presentation.slides)
    shape_count = 0
    chart_count = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            shape_count += 1
            if getattr(shape, "has_chart", False):
                chart_count += 1
    return {"slideCount": slide_count, "shapeCount": shape_count, "chartCount": chart_count}


def _shape_bounds(shape: Any) -> Optional[Tuple[int, int, int, int]]:
    try:
        left = int(getattr(shape, "left"))
        top = int(getattr(shape, "top"))
        width = int(getattr(shape, "width"))
        height = int(getattr(shape, "height"))
    except Exception:
        return None
    if width <= 0 or height <= 0:
        return None
    return left, top, width, height


def _shape_type_name(shape: Any) -> str:
    value = getattr(shape, "shape_type", "")
    return str(getattr(value, "name", value) or "").upper()


def _shape_text_length(shape: Any) -> int:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return 0
    try:
        return len(str(getattr(shape, "text", "") or "").strip())
    except Exception:
        return 0


def _shape_has_visible_payload(shape: Any) -> bool:
    if getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
        return True
    if _shape_type_name(shape) == "PICTURE":
        return True
    return _shape_text_length(shape) > 0


def _is_background_like_shape(
    bounds: Tuple[int, int, int, int],
    slide_width: int,
    slide_height: int,
    slide_area: int,
) -> bool:
    left, top, width, height = bounds
    area = width * height
    return (
        area >= slide_area * 0.85
        and left <= slide_width * 0.03
        and top <= slide_height * 0.03
        and left + width >= slide_width * 0.97
        and top + height >= slide_height * 0.97
    )


def _count_shape_overlaps(bounds: List[Tuple[int, int, int, int, int, str]]) -> int:
    warnings = 0
    for index, first in enumerate(bounds):
        left_a, top_a, width_a, height_a, _shape_index_a, type_a = first
        area_a = width_a * height_a
        for second in bounds[index + 1:]:
            left_b, top_b, width_b, height_b, _shape_index_b, type_b = second
            if type_a == "LINE" or type_b == "LINE":
                continue
            overlap_width = min(left_a + width_a, left_b + width_b) - max(left_a, left_b)
            overlap_height = min(top_a + height_a, top_b + height_b) - max(top_a, top_b)
            if overlap_width <= 0 or overlap_height <= 0:
                continue
            area_b = width_b * height_b
            overlap_area = overlap_width * overlap_height
            if overlap_area / max(1, min(area_a, area_b)) >= PRESENTATION_OVERLAP_RATIO_THRESHOLD:
                warnings += 1
    return warnings


def _presentation_text_role(shape: Any, title_shape: Any, slide_index: int) -> str:
    if title_shape is not None and shape == title_shape:
        return "deck-title" if slide_index == 0 else "slide-title"
    if getattr(shape, "is_placeholder", False):
        try:
            placeholder_name = str(getattr(shape.placeholder_format.type, "name", shape.placeholder_format.type)).upper()
            if "SUBTITLE" in placeholder_name:
                return "subtitle"
        except Exception:
            pass
    return "body"


def _font_check_counts(shape: Any, role: str) -> Tuple[int, int]:
    thresholds = {
        "deck-title": 50.0,
        "slide-title": 35.0,
        "subtitle": 24.0,
        "body": 16.0,
    }
    threshold = thresholds.get(role, 16.0)
    violations = 0
    unspecified = 0
    try:
        paragraphs = list(shape.text_frame.paragraphs)
    except Exception:
        return 0, 0
    for paragraph in paragraphs:
        runs = list(paragraph.runs)
        if not runs and str(getattr(paragraph, "text", "") or "").strip():
            size = getattr(paragraph.font, "size", None)
            if size is None:
                unspecified += 1
            elif float(size.pt) < threshold:
                violations += 1
            continue
        for run in runs:
            if not str(getattr(run, "text", "") or "").strip():
                continue
            size = getattr(run.font, "size", None) or getattr(paragraph.font, "size", None)
            if size is None:
                unspecified += 1
            elif float(size.pt) < threshold:
                violations += 1
    return violations, unspecified


def _chart_has_data(shape: Any) -> bool:
    try:
        chart = shape.chart
        return len(list(chart.series)) > 0
    except Exception:
        return False


def _spreadsheet_chart_has_data(chart: Any) -> bool:
    try:
        return len(list(getattr(chart, "series", []) or [])) > 0
    except Exception:
        return False


def _looks_like_numeric_text(value: str) -> bool:
    text = value.strip()
    if not text or any(char.isalpha() for char in text):
        return False
    return bool(re.fullmatch(r"[$€£¥]?\s*-?\(?\d{1,3}(?:,\d{3})*(?:\.\d+)?\)?%?", text) or re.fullmatch(r"-?\d+(?:\.\d+)?%?", text))


def _looks_like_date_text(value: str) -> bool:
    text = value.strip()
    return bool(
        re.fullmatch(r"\d{4}[-/]\d{1,2}[-/]\d{1,2}", text)
        or re.fullmatch(r"\d{1,2}[-/]\d{1,2}[-/]\d{2,4}", text)
    )


def _length_bucket(length: int) -> str:
    if length <= 0:
        return "empty"
    if length <= 35:
        return "short"
    if length <= 70:
        return "medium"
    return "long"


def _add_check(checks: List[Dict[str, Any]], check_id: str, ok: bool, detail: str = "") -> None:
    checks.append({"id": check_id, "status": "pass" if ok else "fail", "detail": detail})


def _normalize_check(item: Dict[str, Any], gate_names: Iterable[str]) -> Dict[str, Any]:
    status = str(item.get("status") or "").strip().lower()
    if status not in {"pass", "fail", "warn", "pending"}:
        status = "pending"
    allowed_gates = set(gate_names)
    raw_check_id = str(item.get("id") or item.get("gate") or "").strip()
    check_id = raw_check_id if raw_check_id in allowed_gates else "unknown-check"
    return {
        "id": check_id,
        "status": status,
        "detail": _sanitize_check_detail(item.get("detail") or item.get("summary") or ""),
    }


def _sanitize_check_detail(value: Any) -> str:
    allowed_keys = {
        "chart_issues",
        "charts",
        "blank_pages",
        "comment_id_mismatches",
        "comment_refs",
        "comments",
        "date_text",
        "diff",
        "diff_mismatches",
        "empty_sheets",
        "empty_slides",
        "empty_text_pages",
        "error_cells",
        "export",
        "expected_min",
        "extraction_errors",
        "formula_errors",
        "formulas",
        "generated",
        "glyph_issues",
        "image_only_pages",
        "issues",
        "manual_visual_review",
        "missing_titles",
        "non_empty_sheets",
        "numeric_text",
        "out_of_bounds",
        "overlaps",
        "page_count",
        "page_size_variants",
        "pages_compared",
        "paragraphs",
        "rendered",
        "route",
        "rotation_issues",
        "sections",
        "sheets",
        "slides",
        "table_issues",
        "table_candidates",
        "table_text_candidates",
        "tables",
        "text_pages",
        "titles",
        "headings",
        "tracked_changes",
        "unspecified",
        "violations",
    }
    allowed_enums = {
        "artifact-tool",
        "template-following",
        "verified-existing-deck",
        "unspecified",
        "pass",
        "pending",
        "verified",
        "missing",
        "not_applicable",
    }
    parts: List[str] = []
    for raw_part in str(value or "").split(";"):
        if "=" not in raw_part:
            continue
        raw_key, raw_val = raw_part.split("=", 1)
        key = raw_key.strip().lower()
        val = raw_val.strip().lower()
        if key not in allowed_keys:
            continue
        if val.isdigit():
            parts.append(f"{key}={int(val)}")
        elif val in allowed_enums:
            parts.append(f"{key}={val}")
    return "; ".join(parts)[:240]


def _now() -> str:
    return time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
