"""Offline smoke probe for OCR inference and Office format round trips."""

from __future__ import annotations

import io
import importlib.util
import json
from pathlib import Path
import sys


def _activate_pack_runtime(python_root: Path) -> None:
    """Make the staged Pack the only possible third-party import source."""

    if not python_root.is_dir():
        raise RuntimeError("Pack Python runtime is missing")
    retained: list[str] = []
    for raw in sys.path:
        if not raw:
            continue
        parts = {part.casefold() for part in Path(raw).parts}
        if "site-packages" in parts or "dist-packages" in parts:
            continue
        retained.append(raw)
    sys.path[:] = [str(python_root), *retained]


def _module_origins(names: tuple[str, ...], python_root: Path) -> dict[str, str]:
    origins: dict[str, str] = {}
    for name in names:
        module = sys.modules.get(name)
        raw = getattr(module, "__file__", None)
        if not isinstance(raw, str):
            raise RuntimeError(f"Module origin is unavailable: {name}")
        try:
            relative = Path(raw).resolve(strict=True).relative_to(python_root)
        except (OSError, ValueError):
            raise RuntimeError(f"Module escaped Pack runtime: {name}") from None
        origins[name] = relative.as_posix()
    return origins


def _office(worker_path: Path, python_root: Path) -> dict[str, object]:
    from docx import Document
    from openpyxl import Workbook, load_workbook
    from pptx import Presentation
    from pypdf import PdfReader
    from reportlab.pdfgen import canvas

    spec = importlib.util.spec_from_file_location("ecorex_dependency_pack_worker", worker_path)
    if spec is None or spec.loader is None:
        raise RuntimeError("Office worker cannot be loaded")
    worker = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(worker)

    document_bytes = io.BytesIO()
    document = Document()
    document.add_paragraph("EcoreX Office Pack")
    document.save(document_bytes)

    workbook_bytes = io.BytesIO()
    workbook = Workbook()
    workbook.active["A1"] = "EcoreX"
    workbook.save(workbook_bytes)

    presentation_bytes = io.BytesIO()
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(presentation_bytes)

    pdf_bytes = io.BytesIO()
    pdf = canvas.Canvas(pdf_bytes)
    pdf.drawString(32, 780, "EcoreX PDF")
    pdf.save()
    page_count = len(PdfReader(io.BytesIO(pdf_bytes.getvalue())).pages)
    document_text = Document(io.BytesIO(document_bytes.getvalue())).paragraphs[0].text
    spreadsheet_text = load_workbook(
        io.BytesIO(workbook_bytes.getvalue()), read_only=True
    ).active["A1"].value
    presentation_pages = len(
        Presentation(io.BytesIO(presentation_bytes.getvalue())).slides
    )
    if (
        document_text != "EcoreX Office Pack"
        or spreadsheet_text != "EcoreX"
        or presentation_pages != 1
        or page_count != 1
    ):
        raise RuntimeError("Office round-trip contract is invalid")
    authoring_cases = {
        "document": {
            "title": "v1",
            "sections": [{"heading": "Summary", "paragraphs": ["v1"]}],
        },
        "spreadsheet": {
            "title": "v1",
            "sheets": [{"name": "Data", "rows": [["version", 1]]}],
        },
        "presentation": {
            "title": "v1",
            "slides": [{"title": "Summary", "bullets": ["v1"]}],
        },
        "pdf": {
            "title": "v1",
            "sections": [{"heading": "Summary", "paragraphs": ["v1"]}],
        },
    }
    edited = {}
    for family, payload in authoring_cases.items():
        created = worker._office_create({"family": family, **payload}, python_root)
        edited_payload = {
            "document": {
                "title": "v2",
                "sections": [{"heading": "Summary", "paragraphs": ["v2"]}],
            },
            "spreadsheet": {
                "title": "v2",
                "sheets": [{"name": "Data", "rows": [["version", "v2"]]}],
            },
            "presentation": {
                "title": "v2",
                "slides": [{"title": "Summary v2", "bullets": ["v2"]}],
            },
            "pdf": {
                "title": "v2",
                "sections": [{"heading": "Summary", "paragraphs": ["v2"]}],
            },
        }[family]
        updated = worker._office_edit(
            {
                "family": family,
                **edited_payload,
                "content_base64": created["content_base64"],
            },
            python_root,
        )
        inspection = worker._office_read(
            {
                "family": family,
                "content_base64": updated["content_base64"],
            },
            python_root,
        )
        if (
            updated.get("validation", {}).get("source_opened") is not True
            or "v2" not in inspection.get("text", "")
        ):
            raise RuntimeError("Office public worker edit contract is invalid")
        edited[family] = True
    return {
        "document_bytes": len(document_bytes.getvalue()),
        "spreadsheet_bytes": len(workbook_bytes.getvalue()),
        "presentation_bytes": len(presentation_bytes.getvalue()),
        "pdf_bytes": len(pdf_bytes.getvalue()),
        "pdf_pages": page_count,
        "round_trip": True,
        "worker_create_edit": edited,
    }


def _ocr() -> dict[str, object]:
    import numpy
    import cv2  # noqa: F401
    import onnxruntime  # noqa: F401
    import pyclipper  # noqa: F401
    from PIL import Image, ImageDraw, ImageFont
    from rapidocr_onnxruntime import RapidOCR

    image = Image.new("RGB", (640, 180), "white")
    font = ImageFont.load_default(size=48)
    ImageDraw.Draw(image).text((30, 50), "ECOREX 2026", font=font, fill="black")
    result, elapsed = RapidOCR()(numpy.asarray(image))
    if result is not None and not isinstance(result, list):
        raise RuntimeError("OCR result contract is invalid")
    texts = [str(item[1]) for item in result or () if len(item) >= 2]
    normalized = "".join(character for text in texts for character in text if character.isalnum()).upper()
    if "ECOREX2026" not in normalized:
        raise RuntimeError("OCR model did not recognize the offline fixture")
    return {
        "result_count": len(result or ()),
        "elapsed_reported": elapsed is not None,
        "image_size": list(image.size),
        "fixture_recognized": True,
    }


def main() -> int:
    if len(sys.argv) not in {3, 4} or sys.argv[1] not in {"ocr", "office"}:
        return 64
    root = Path(sys.argv[2]).resolve(strict=True)
    python_root = (root / "runtime" / "python").resolve(strict=True)
    _activate_pack_runtime(python_root)
    pack_id = sys.argv[1]
    if pack_id == "office" and len(sys.argv) != 4:
        return 64
    result = (
        _ocr()
        if pack_id == "ocr"
        else _office(Path(sys.argv[3]).resolve(strict=True), python_root)
    )
    required_modules = (
        ("rapidocr_onnxruntime", "onnxruntime", "numpy", "PIL", "cv2", "pyclipper")
        if pack_id == "ocr"
        else ("docx", "openpyxl", "pptx", "pypdf", "reportlab")
    )
    origins = _module_origins(required_modules, python_root)
    print(
        json.dumps(
            {
                "schema_version": 1,
                "pack_id": pack_id,
                "isolation": {
                    "mode": "pack-only-third-party",
                    "module_origins": origins,
                },
                "result": result,
            },
            sort_keys=True,
            separators=(",", ":"),
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
