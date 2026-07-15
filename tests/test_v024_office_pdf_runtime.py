from __future__ import annotations

import json
import inspect

import pytest

from common.office_pdf_runtime import (
    ARTIFACT_KINDS,
    _hash_ref,
    _render_artifact_proof,
    _register_trusted_render_artifact,
    analyze_document_quality,
    analyze_pdf_quality,
    analyze_spreadsheet_quality,
    analyze_presentation_quality,
    build_document_quality_evidence,
    build_pdf_quality_evidence,
    build_spreadsheet_quality_evidence,
    build_presentation_quality_evidence,
    build_quality_evidence,
    compare_pdf_page_quality,
    default_quality_gates,
    detect_artifact_kind,
    inspect_office_pdf_artifact,
    probe_office_pdf_runtime,
    render_presentation_preview,
    render_document_preview,
    render_spreadsheet_preview,
    render_pdf_pages,
)


def test_office_pdf_runtime_detects_supported_kinds():
    assert detect_artifact_kind("deck.pptx") == "presentation"
    assert detect_artifact_kind("model.xlsx") == "spreadsheet"
    assert detect_artifact_kind("model.tsv") == "spreadsheet"
    assert detect_artifact_kind("legacy.xls") == "unknown"
    assert detect_artifact_kind("notes.docx") == "document"
    assert detect_artifact_kind("report.pdf") == "pdf"
    assert detect_artifact_kind("unknown.bin") == "unknown"


def test_office_pdf_runtime_probe_is_redacted_and_maps_quality_gates():
    payload = probe_office_pdf_runtime()
    serialized = json.dumps(payload, ensure_ascii=False)

    assert payload["schemaVersion"] == 1
    assert payload["packId"] == "office-pdf"
    assert payload["redacted"] is True
    assert "python.exe" not in serialized
    assert "C:\\" not in serialized
    for kind, spec in ARTIFACT_KINDS.items():
        row = payload["artifactKinds"][kind]
        assert row["compatibilityId"] == spec["compatibilityId"]
        assert row["officialSkill"] == spec["officialSkill"]
        assert row["qualityGates"] == spec["qualityGates"]


def test_quality_evidence_fails_when_required_gate_missing(tmp_path):
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"placeholder")
    gates = default_quality_gates("presentation")
    checks = [{"id": gate, "status": "pass"} for gate in gates[:-1]]

    payload = build_quality_evidence("presentation", checks, source=source)
    serialized = json.dumps(payload, ensure_ascii=False)

    assert payload["status"] == "fail"
    assert payload["missingQualityGates"] == [gates[-1]]
    assert str(tmp_path) not in serialized
    assert source.name not in serialized


def test_quality_evidence_sanitizes_untrusted_render_metadata(tmp_path):
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"placeholder")
    gates = default_quality_gates("presentation")
    checks = [{"id": gate, "status": "pass"} for gate in gates]
    renders = [
        {
            "slide": 1,
            "artifactRef": r"C:\private\private-slide-001.png",
            "sourceRef": "https://internal.example/render/private-slide-001.png",
            "extension": ".png",
            "sizeBytes": 2048,
            "width": 1280,
            "height": 720,
            "path": r"C:\private\private-slide-001.png",
            "fileName": "private-slide-001.png",
            "content": "Private render notes",
            "text": "Secret visible caption",
            "url": "https://internal.example/render/private-slide-001.png",
        }
    ]

    payload = build_quality_evidence("presentation", checks, source=source, renders=renders)
    serialized = json.dumps(payload, ensure_ascii=False)
    artifact = payload["renderedArtifacts"][0]

    assert set(artifact) == {"slide", "sizeBytes", "width", "height", "extension", "artifactRef", "sourceRef"}
    assert artifact["artifactRef"].startswith("hmac:")
    assert artifact["sourceRef"].startswith("hmac:")
    assert "private-slide-001.png" not in serialized
    assert "Private render notes" not in serialized
    assert "Secret visible caption" not in serialized
    assert "internal.example" not in serialized
    assert "C:\\" not in serialized


def test_quality_evidence_sanitizes_untrusted_check_id_and_detail(tmp_path):
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"placeholder")
    gates = default_quality_gates("presentation")
    checks = [
        {
            "id": r"story-flow C:\private\customer-deck.pptx",
            "status": "pass",
            "detail": r"Private check detail C:\private\customer-deck.pptx Secret content; rendered=1",
        },
        *[{"id": gate, "status": "pass"} for gate in gates if gate != "story-flow"],
    ]

    payload = build_quality_evidence("presentation", checks, source=source)
    serialized = json.dumps(payload, ensure_ascii=False)
    unknown = next(check for check in payload["checks"] if check["id"] == "unknown-check")

    assert payload["status"] == "fail"
    assert "story-flow" in payload["missingQualityGates"]
    assert unknown["detail"] == "rendered=1"
    assert "customer-deck.pptx" not in serialized
    assert "Private check detail" not in serialized
    assert "Secret content" not in serialized
    assert "C:\\" not in serialized


def test_pdf_render_api_does_not_expose_raw_path_option():
    assert "include_paths" not in inspect.signature(render_pdf_pages).parameters
    assert "include_paths" not in inspect.signature(render_presentation_preview).parameters
    assert "include_paths" not in inspect.signature(render_spreadsheet_preview).parameters
    assert "include_paths" not in inspect.signature(render_document_preview).parameters


def _create_private_pdf(path, *, blank=False, landscape_page=False, second_page=False):
    pytest.importorskip("reportlab")
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas

    page_size = landscape(letter) if landscape_page else letter
    c = canvas.Canvas(str(path), pagesize=page_size)
    if not blank:
        c.setFont("Helvetica", 14)
        c.drawString(72, page_size[1] - 72, "Private PDF report")
        c.setFont("Helvetica", 10)
        rows = [
            ("Private Product", "Units", "Price", "Revenue"),
            ("Secret Alpha", "10", "2.5", "25"),
            ("Secret Beta", "12", "3.0", "36"),
        ]
        start_x = 72
        start_y = page_size[1] - 120
        col_width = 95
        row_height = 22
        for row_index, row in enumerate(rows):
            y = start_y - row_index * row_height
            for col_index, value in enumerate(row):
                x = start_x + col_index * col_width
                c.drawString(x + 4, y - 15, value)
                c.rect(x, y - row_height, col_width, row_height)
    c.showPage()
    if second_page:
        c.setFont("Helvetica", 12)
        c.drawString(72, page_size[1] - 72, "Private appendix")
        c.showPage()
    c.save()


def _trusted_pdf_render(source, page=1, width=612, height=792):
    item = {
        "page": page,
        "artifactRef": f"hmac:{page:016d}",
        "extension": ".png",
        "sizeBytes": 4096,
        "width": width,
        "height": height,
    }
    item["renderProof"] = _register_trusted_render_artifact(_hash_ref(source), item)
    return item


def test_pdf_quality_evidence_passes_clean_pdf_without_text_leaks(tmp_path):
    path = tmp_path / "clean-private-report.pdf"
    _create_private_pdf(path)

    renders = [_trusted_pdf_render(path)]
    evidence = build_pdf_quality_evidence(path, renders=renders, reference_path=path, visual_inspection_passed=True)
    analysis = analyze_pdf_quality(path)
    diff = compare_pdf_page_quality(path, path)
    serialized = json.dumps({"evidence": evidence, "analysis": analysis, "diff": diff}, ensure_ascii=False)

    assert evidence["status"] == "pass"
    assert evidence["missingQualityGates"] == []
    assert evidence["pdfAnalysis"]["summary"]["pageCount"] == 1
    assert evidence["pdfAnalysis"]["summary"]["textExtractablePageCount"] == 1
    assert evidence["pdfAnalysis"]["summary"]["tableCandidatePageCount"] >= 1
    assert evidence["pdfDiffAnalysis"]["summary"]["mismatchCount"] == 0
    assert diff["summary"]["mismatchCount"] == 0
    assert "Private PDF report" not in serialized
    assert "Secret Alpha" not in serialized
    assert path.name not in serialized


def test_pdf_quality_evidence_detects_blank_pages_and_rejects_fake_render(tmp_path):
    path = tmp_path / "blank-private-report.pdf"
    _create_private_pdf(path, blank=True)

    evidence = build_pdf_quality_evidence(path, renders=[{"page": 1, "artifactRef": "hmac:deadbeef"}])
    serialized = json.dumps(evidence, ensure_ascii=False)
    failed_checks = {check["id"] for check in evidence["checks"] if check["status"] == "fail"}

    assert evidence["status"] == "fail"
    assert {"text-orientation", "page-render", "layout-inspection"} <= failed_checks
    assert evidence["renderedArtifacts"] == []
    assert evidence["pdfAnalysis"]["summary"]["blankPageRiskCount"] >= 1
    assert path.name not in serialized


def test_pdf_quality_evidence_rejects_forged_render_proof_without_runtime_registration(tmp_path):
    path = tmp_path / "forged-private-report.pdf"
    _create_private_pdf(path)
    forged = {
        "page": 1,
        "artifactRef": "hmac:deadbeefdeadbeef",
        "extension": ".png",
        "sizeBytes": 4096,
        "width": 612,
        "height": 792,
    }
    forged["renderProof"] = _render_artifact_proof(forged)

    evidence = build_pdf_quality_evidence(path, renders=[forged])
    page_render = next(check for check in evidence["checks"] if check["id"] == "page-render")

    assert evidence["status"] == "fail"
    assert evidence["renderedArtifacts"] == []
    assert page_render["status"] == "fail"
    assert "rendered=0" in page_render["detail"]


def test_pdf_visual_diff_detects_reference_layout_mismatch_without_leaks(tmp_path):
    reference = tmp_path / "reference-private-report.pdf"
    candidate = tmp_path / "candidate-private-report.pdf"
    _create_private_pdf(reference)
    _create_private_pdf(candidate, landscape_page=True)

    renders = [_trusted_pdf_render(candidate, width=792, height=612)]
    evidence = build_pdf_quality_evidence(candidate, renders=renders, reference_path=reference)
    serialized = json.dumps(evidence, ensure_ascii=False)
    visual_diff = next(check for check in evidence["checks"] if check["id"] == "visual-diff")

    assert evidence["status"] == "fail"
    assert visual_diff["status"] == "fail"
    assert evidence["pdfDiffAnalysis"]["summary"]["mismatchCount"] >= 1
    assert evidence["pdfDiffAnalysis"]["summary"]["orientationMismatchCount"] >= 1
    assert "Private PDF report" not in serialized
    assert "Secret Alpha" not in serialized
    assert reference.name not in serialized
    assert candidate.name not in serialized


def test_docx_inspection_reports_counts_without_text(tmp_path):
    docx = pytest.importorskip("docx")
    document = docx.Document()
    document.add_paragraph("secret customer paragraph")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "hidden cell"
    path = tmp_path / "private-notes.docx"
    document.save(path)

    payload = inspect_office_pdf_artifact(path)
    serialized = json.dumps(payload, ensure_ascii=False)

    assert payload["kind"] == "document"
    assert payload["summary"]["paragraphCount"] >= 1
    assert payload["summary"]["tableCount"] == 1
    assert "secret customer paragraph" not in serialized
    assert "hidden cell" not in serialized
    assert path.name not in serialized


def _break_docx_table_geometry_and_add_orphan_comment(path):
    import zipfile
    import xml.etree.ElementTree as ET

    namespace = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    ET.register_namespace("w", namespace)
    with zipfile.ZipFile(path, "r") as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}
    root = ET.fromstring(entries["word/document.xml"])
    for table in root.findall(f".//{{{namespace}}}tbl"):
        for child in list(table):
            if child.tag == f"{{{namespace}}}tblGrid":
                table.remove(child)
        properties = table.find(f"{{{namespace}}}tblPr")
        if properties is not None:
            for child in list(properties):
                if child.tag == f"{{{namespace}}}tblW":
                    properties.remove(child)
        for cell in table.findall(f".//{{{namespace}}}tc"):
            cell_properties = cell.find(f"{{{namespace}}}tcPr")
            if cell_properties is not None:
                for child in list(cell_properties):
                    if child.tag == f"{{{namespace}}}tcW":
                        cell_properties.remove(child)
    entries["word/document.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    entries["word/comments.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        b'<w:comment w:id="0" w:author="Private reviewer" w:date="2026-01-01T00:00:00Z">'
        b"<w:p><w:r><w:t>Secret orphan comment</w:t></w:r></w:p>"
        b"</w:comment></w:comments>"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)


def _add_mismatched_docx_comment_reference(path):
    import zipfile
    import xml.etree.ElementTree as ET

    namespace = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    ET.register_namespace("w", namespace)
    with zipfile.ZipFile(path, "r") as archive:
        entries = {name: archive.read(name) for name in archive.namelist()}
    root = ET.fromstring(entries["word/document.xml"])
    paragraph = root.find(f".//{{{namespace}}}p")
    if paragraph is not None:
        run = ET.SubElement(paragraph, f"{{{namespace}}}r")
        ET.SubElement(run, f"{{{namespace}}}commentReference", {f"{{{namespace}}}id": "1"})
    entries["word/document.xml"] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
    entries["word/comments.xml"] = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        b'<w:comment w:id="0" w:author="Private reviewer" w:date="2026-01-01T00:00:00Z">'
        b"<w:p><w:r><w:t>Secret mismatched comment</w:t></w:r></w:p>"
        b"</w:comment></w:comments>"
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)


def test_document_quality_evidence_passes_clean_docx_without_text_leaks(tmp_path):
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Secret document title", 0)
    document.add_heading("Private section", level=1)
    document.add_paragraph("Confidential customer paragraph")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Secret header"
    table.cell(0, 1).text = "Private header"
    table.cell(1, 0).text = "Hidden value"
    table.cell(1, 1).text = "Internal value"
    path = tmp_path / "clean-secret-doc.docx"
    document.save(path)

    renders = [{"page": 1, "artifactRef": "hmac:3333333333333333", "width": 1280, "height": 720}]
    for item in renders:
        item["renderProof"] = _render_artifact_proof(item)

    evidence = build_document_quality_evidence(path, renders=renders, visual_inspection_passed=True)
    analysis = analyze_document_quality(path)
    serialized = json.dumps({"evidence": evidence, "analysis": analysis}, ensure_ascii=False)

    assert evidence["status"] == "pass"
    assert evidence["documentAnalysis"]["summary"]["tableCount"] == 1
    assert evidence["documentAnalysis"]["summary"]["tableIssueCount"] == 0
    assert evidence["documentAnalysis"]["summary"]["headingParagraphCount"] >= 1
    assert "Secret document title" not in serialized
    assert "Confidential customer paragraph" not in serialized
    assert "Hidden value" not in serialized
    assert path.name not in serialized


def test_document_quality_evidence_detects_structure_table_redline_and_render_failures(tmp_path):
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_paragraph("Secret body without heading")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Private cell"
    path = tmp_path / "bad-secret-doc.docx"
    document.save(path)
    _break_docx_table_geometry_and_add_orphan_comment(path)

    evidence = build_document_quality_evidence(
        path,
        renders=[{"page": 1, "artifactRef": "hmac:deadbeef"}],
        visual_inspection_passed=True,
    )
    serialized = json.dumps(evidence, ensure_ascii=False)
    failed_checks = {check["id"] for check in evidence["checks"] if check["status"] == "fail"}

    assert evidence["status"] == "fail"
    assert {"design-preset", "render-docx", "table-geometry", "redline-preserve"} <= failed_checks
    assert evidence["renderedArtifacts"] == []
    assert evidence["documentAnalysis"]["summary"]["tableIssueCount"] >= 1
    assert evidence["documentAnalysis"]["summary"]["commentCount"] == 1
    assert evidence["documentAnalysis"]["summary"]["commentReferenceCount"] == 0
    assert "Secret body without heading" not in serialized
    assert "Private cell" not in serialized
    assert "Secret orphan comment" not in serialized
    assert path.name not in serialized


def test_document_quality_evidence_detects_equal_count_comment_id_mismatch(tmp_path):
    docx = pytest.importorskip("docx")

    document = docx.Document()
    document.add_heading("Secret comment mismatch title", 0)
    document.add_paragraph("Private body")
    path = tmp_path / "comment-mismatch-secret-doc.docx"
    document.save(path)
    _add_mismatched_docx_comment_reference(path)

    renders = [{"page": 1, "artifactRef": "hmac:4444444444444444", "width": 1280, "height": 720}]
    for item in renders:
        item["renderProof"] = _render_artifact_proof(item)

    evidence = build_document_quality_evidence(path, renders=renders, visual_inspection_passed=True)
    serialized = json.dumps(evidence, ensure_ascii=False)
    redline_check = next(check for check in evidence["checks"] if check["id"] == "redline-preserve")

    assert evidence["status"] == "fail"
    assert redline_check["status"] == "fail"
    assert evidence["documentAnalysis"]["summary"]["commentCount"] == 1
    assert evidence["documentAnalysis"]["summary"]["commentReferenceCount"] == 1
    assert evidence["documentAnalysis"]["summary"]["commentIdMismatchCount"] == 2
    assert "Secret mismatched comment" not in serialized
    assert "Secret comment mismatch title" not in serialized
    assert path.name not in serialized


def test_xlsx_inspection_reports_structure_without_cell_values(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A1"] = "secret customer value"
    sheet["B1"] = "=SUM(1,2)"
    path = tmp_path / "private-model.xlsx"
    workbook.save(path)

    payload = inspect_office_pdf_artifact(path)
    serialized = json.dumps(payload, ensure_ascii=False)

    assert payload["kind"] == "spreadsheet"
    assert payload["summary"]["sheetCount"] == 1
    assert payload["summary"]["formulaCellCount"] == 1
    assert payload["summary"]["cellsScanned"] >= 2
    assert payload["summary"]["truncated"] is False
    assert "secret customer value" not in serialized
    assert path.name not in serialized


def test_spreadsheet_quality_evidence_passes_clean_workbook_without_text_leaks(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.chart import BarChart, Reference

    workbook = openpyxl.Workbook()
    inputs = workbook.active
    inputs.title = "Private Inputs"
    inputs.append(["Product", "Units", "Price", "Revenue"])
    inputs.append(["Secret Alpha", 10, 2.5, "=B2*C2"])
    inputs.append(["Secret Beta", 12, 3.0, "=B3*C3"])
    inputs.append(["Secret Gamma", 8, 4.5, "=B4*C4"])
    dashboard = workbook.create_sheet("Private Dashboard")
    dashboard["A1"] = "Secret revenue dashboard"
    dashboard["B2"] = "Total"
    dashboard["C2"] = "=SUM('Private Inputs'!D2:D4)"
    chart = BarChart()
    chart.title = "Revenue by product"
    data = Reference(inputs, min_col=4, min_row=1, max_row=4)
    categories = Reference(inputs, min_col=1, min_row=2, max_row=4)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    dashboard.add_chart(chart, "E2")
    path = tmp_path / "clean-secret-workbook.xlsx"
    workbook.save(path)

    renders = [
        {"page": 1, "artifactRef": "hmac:1111111111111111", "width": 1280, "height": 720},
        {"page": 2, "artifactRef": "hmac:2222222222222222", "width": 1280, "height": 720},
    ]
    for item in renders:
        item["renderProof"] = _render_artifact_proof(item)

    evidence = build_spreadsheet_quality_evidence(path, renders=renders, visual_inspection_passed=True)
    analysis = analyze_spreadsheet_quality(path)
    serialized = json.dumps({"evidence": evidence, "analysis": analysis}, ensure_ascii=False)

    assert evidence["status"] == "pass"
    assert evidence["spreadsheetAnalysis"]["summary"]["chartCount"] == 1
    assert evidence["spreadsheetAnalysis"]["summary"]["formulaCellCount"] >= 4
    assert evidence["spreadsheetAnalysis"]["summary"]["numericTextRiskCount"] == 0
    assert "Secret Alpha" not in serialized
    assert "Secret revenue dashboard" not in serialized
    assert "Private Inputs" not in serialized
    assert path.name not in serialized


def test_spreadsheet_quality_evidence_detects_typed_formula_chart_and_render_failures(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    from openpyxl.chart import BarChart

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Private Bad Sheet"
    sheet["A1"] = "100"
    sheet["B1"] = "2026-01-01"
    sheet["C1"] = "=#REF!"
    sheet["D1"] = "Secret bad workbook"
    sheet.add_chart(BarChart(), "F2")
    path = tmp_path / "bad-secret-workbook.xlsx"
    workbook.save(path)

    evidence = build_spreadsheet_quality_evidence(
        path,
        renders=[{"page": 1, "artifactRef": "hmac:deadbeef"}],
        visual_inspection_passed=True,
    )
    serialized = json.dumps(evidence, ensure_ascii=False)
    failed_checks = {check["id"] for check in evidence["checks"] if check["status"] == "fail"}

    assert evidence["status"] == "fail"
    assert {"typed-values", "formula-audit", "chart-render", "render-preview"} <= failed_checks
    assert evidence["renderedArtifacts"] == []
    assert evidence["spreadsheetAnalysis"]["summary"]["numericTextRiskCount"] >= 1
    assert evidence["spreadsheetAnalysis"]["summary"]["dateTextRiskCount"] >= 1
    assert evidence["spreadsheetAnalysis"]["summary"]["formulaErrorTokenCount"] >= 1
    assert evidence["spreadsheetAnalysis"]["summary"]["chartIssueCount"] >= 1
    assert "Secret bad workbook" not in serialized
    assert "Private Bad Sheet" not in serialized
    assert "#REF!" not in serialized
    assert path.name not in serialized


def test_pptx_inspection_reports_structure_without_text(tmp_path):
    pptx = pytest.importorskip("pptx")
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "secret customer title"
    path = tmp_path / "private-deck.pptx"
    presentation.save(path)

    payload = inspect_office_pdf_artifact(path)
    serialized = json.dumps(payload, ensure_ascii=False)

    assert payload["kind"] == "presentation"
    assert payload["summary"]["slideCount"] == 1
    assert payload["summary"]["shapeCount"] >= 1
    assert "secret customer title" not in serialized
    assert path.name not in serialized


def test_presentation_quality_evidence_passes_clean_deck_without_text_leaks(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = title_slide.shapes.title
    title.text = "Confidential launch narrative"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(54)
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Private customer plan"
    subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_slide.shapes.title.text = "Revenue momentum"
    chart_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(38)
    body = chart_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.2), Inches(1.4))
    body.text_frame.text = "Secret metric summary"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Series 1", (1, 2, 3))
    chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.1),
        Inches(1.4),
        Inches(4.2),
        Inches(3.9),
        chart_data,
    )

    path = tmp_path / "clean-secret-deck.pptx"
    presentation.save(path)
    renders = [
        {"slide": 1, "artifactRef": "hmac:0123456789abcdef", "width": 1280, "height": 720},
        {"slide": 2, "artifactRef": "hmac:fedcba9876543210", "width": 1280, "height": 720},
    ]
    for item in renders:
        item["renderProof"] = _render_artifact_proof(item)

    evidence = build_presentation_quality_evidence(
        path,
        authoring_route="artifact-tool",
        renders=renders,
        visual_inspection_passed=True,
    )
    serialized = json.dumps(evidence, ensure_ascii=False)

    assert evidence["status"] == "pass"
    assert evidence["missingQualityGates"] == []
    assert evidence["presentationAnalysis"]["summary"]["chartCount"] == 1
    assert evidence["presentationAnalysis"]["summary"]["overlapWarningCount"] == 0
    assert "Confidential launch narrative" not in serialized
    assert "Private customer plan" not in serialized
    assert "Secret metric summary" not in serialized
    assert path.name not in serialized


def test_presentation_quality_evidence_sanitizes_untrusted_authoring_route(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Pt

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Private title"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(38)
    path = tmp_path / "route-secret-deck.pptx"
    presentation.save(path)

    evidence = build_presentation_quality_evidence(
        path,
        authoring_route=r"artifact-tool C:\private\customer-route.pptx Secret customer route",
        renders=[],
        visual_inspection_passed=True,
    )
    serialized = json.dumps(evidence, ensure_ascii=False)

    assert evidence["authoringRoute"] == "unspecified"
    assert any(check["id"] == "artifact-tool-authoring" and check["detail"] == "route=unspecified" for check in evidence["checks"])
    assert "customer-route.pptx" not in serialized
    assert "Secret customer route" not in serialized
    assert "C:\\" not in serialized


def test_presentation_quality_evidence_counts_only_render_refs(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Pt

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Private title"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(38)
    path = tmp_path / "render-ref-deck.pptx"
    presentation.save(path)

    path_only = build_presentation_quality_evidence(
        path,
        authoring_route="artifact-tool",
        renders=[{"path": r"C:\private\fake.png"}],
        visual_inspection_passed=True,
    )
    slide_only = build_presentation_quality_evidence(
        path,
        authoring_route="artifact-tool",
        renders=[{"slide": 1}],
        visual_inspection_passed=True,
    )
    ref_without_proof = build_presentation_quality_evidence(
        path,
        authoring_route="artifact-tool",
        renders=[{"slide": 1, "artifactRef": "hmac:deadbeef"}],
        visual_inspection_passed=True,
    )

    for evidence in (path_only, slide_only, ref_without_proof):
        render_check = next(check for check in evidence["checks"] if check["id"] == "render-preview")
        assert evidence["status"] == "fail"
        assert evidence["renderedArtifacts"] == []
        assert render_check["status"] == "fail"
        assert "rendered=0" in render_check["detail"]


def test_presentation_quality_evidence_detects_layout_and_render_failures(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Tiny title"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    first = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(3.0), Inches(1.2))
    first.text_frame.text = "hidden overlap one"
    first.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    second = slide.shapes.add_textbox(Inches(1.4), Inches(1.7), Inches(3.0), Inches(1.2))
    second.text_frame.text = "hidden overlap two"
    second.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    out_of_bounds = slide.shapes.add_textbox(Inches(12.8), Inches(7.2), Inches(1.0), Inches(1.0))
    out_of_bounds.text = "hidden clipped"

    path = tmp_path / "bad-secret-deck.pptx"
    presentation.save(path)

    evidence = build_presentation_quality_evidence(path, authoring_route="", renders=[], visual_inspection_passed=False)
    analysis = analyze_presentation_quality(path)
    serialized = json.dumps({"evidence": evidence, "analysis": analysis}, ensure_ascii=False)

    assert evidence["status"] == "fail"
    assert evidence["presentationAnalysis"]["summary"]["fontViolationCount"] >= 1
    assert evidence["presentationAnalysis"]["summary"]["overlapWarningCount"] >= 1
    assert evidence["presentationAnalysis"]["summary"]["outOfBoundsCount"] >= 1
    failed_checks = {check["id"] for check in evidence["checks"] if check["status"] == "fail"}
    assert {"layout-bounds", "font-size-check", "render-preview", "overlap-check"} <= failed_checks
    assert "hidden overlap" not in serialized
    assert path.name not in serialized
