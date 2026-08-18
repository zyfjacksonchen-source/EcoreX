"""Signed Pack-Python worker for dependency service invocations."""

from __future__ import annotations

import base64
from io import BytesIO
import json
import os
from pathlib import Path, PurePosixPath
import subprocess
import sys
import textwrap
import time
from typing import Any, Mapping
import zipfile


_MAX_OFFICE_CONTENT_BYTES = 5 * 1024 * 1024
_MAX_OFFICE_TEXT_BYTES = 192 * 1024
_MAX_OOXML_FILES = 4096
_MAX_OOXML_UNCOMPRESSED_BYTES = 64 * 1024 * 1024
# ponytail: one measured Office worker cap keeps the contract small; split it
# per family only if a real packaged document proves 512 MiB insufficient.
OFFICE_READ_PROCESS_MEMORY_LIMIT_BYTES = 512 * 1024 * 1024
OFFICE_READ_JOB_MEMORY_LIMIT_BYTES = 768 * 1024 * 1024


def _inside(module: Any, root: Path) -> None:
    origin = getattr(module, "__file__", None)
    try:
        Path(origin).resolve(strict=True).relative_to(root.resolve(strict=True))
    except (OSError, TypeError, ValueError):
        raise RuntimeError("dependency origin escaped verified pack") from None


def _ocr(payload: Mapping[str, Any], runtime: Path) -> Mapping[str, Any]:
    import numpy as np
    from PIL import Image, ImageOps
    import PIL
    import rapidocr_onnxruntime
    from rapidocr_onnxruntime import RapidOCR

    for module in (np, PIL, rapidocr_onnxruntime):
        _inside(module, runtime)
    content = base64.b64decode(str(payload.get("content_base64") or ""), validate=True)
    if not content or len(content) > 8 * 1024 * 1024:
        raise ValueError("OCR content is invalid")
    started = time.monotonic()
    with Image.open(BytesIO(content)) as source:
        pixels = np.asarray(ImageOps.exif_transpose(source).convert("RGB"))
    raw = RapidOCR()(pixels)
    values = raw[0] if isinstance(raw, tuple) and raw else raw
    blocks: list[dict[str, Any]] = []
    if isinstance(values, list):
        for item in values:
            if not isinstance(item, (list, tuple)) or len(item) < 2:
                continue
            text = str(item[1] or "").strip()
            if not text:
                continue
            confidence = item[2] if len(item) > 2 else None
            blocks.append(
                {
                    "text": text,
                    "confidence": (
                        max(0.0, min(1.0, float(confidence)))
                        if isinstance(confidence, (int, float))
                        and not isinstance(confidence, bool)
                        else None
                    ),
                }
            )
    text = "\n".join(item["text"] for item in blocks)
    return {
        "status": "success" if text else "empty",
        "provider": "rapidocr_onnxruntime",
        "text": text,
        "blocks": blocks,
        "latencyMs": int((time.monotonic() - started) * 1000),
        "cacheHit": False,
    }


def _office_probe(runtime: Path) -> Mapping[str, Any]:
    import docx
    import openpyxl
    import pptx
    import pypdf
    import reportlab

    for module in (docx, openpyxl, pptx, pypdf, reportlab):
        _inside(module, runtime)
    return {
        "provider": "python-office-formats-v1",
        "modules": ["docx", "openpyxl", "pptx", "pypdf", "reportlab"],
    }


def _text(value: Any, *, maximum: int = 4096) -> str:
    text = str(value or "").strip()
    if not text or len(text.encode("utf-8")) > maximum:
        raise ValueError("office text is invalid")
    return text


class _TextCollector:
    def __init__(self) -> None:
        self.parts: list[str] = []
        self.size_bytes = 0
        self.truncated = False

    def add(self, value: Any) -> bool:
        text = str(value or "").strip()
        if not text:
            return True
        encoded = (text + "\n").encode("utf-8")
        remaining = _MAX_OFFICE_TEXT_BYTES - self.size_bytes
        if len(encoded) > remaining:
            if remaining > 0:
                self.parts.append(encoded[:remaining].decode("utf-8", errors="ignore"))
                self.size_bytes = _MAX_OFFICE_TEXT_BYTES
            self.truncated = True
            return False
        self.parts.append(text + "\n")
        self.size_bytes += len(encoded)
        return True

    def value(self) -> str:
        return "".join(self.parts).rstrip()


def _office_content(payload: Mapping[str, Any]) -> tuple[str, bytes]:
    family = str(payload.get("family") or "")
    if family not in {"document", "spreadsheet", "presentation", "pdf"}:
        raise ValueError("office family is unsupported")
    content = base64.b64decode(str(payload.get("content_base64") or ""), validate=True)
    if not 1 <= len(content) <= _MAX_OFFICE_CONTENT_BYTES:
        raise ValueError("office input size is invalid")
    return family, content


def _apply_office_read_memory_limit(limit_bytes: int) -> None:
    if limit_bytes != OFFICE_READ_PROCESS_MEMORY_LIMIT_BYTES:
        raise RuntimeError("office read memory limit is invalid")
    if os.name == "nt":
        # The parent starts Windows workers suspended and assigns the hard
        # Process/Job limits before resuming them.
        return
    try:
        import resource

        inherited_hard = resource.getrlimit(resource.RLIMIT_AS)[1]
        absolute_limit = limit_bytes
        if sys.platform == "darwin":
            # arm64 macOS maps a very large shared address region before Python
            # starts. Bound new address space above that measured baseline.
            measured = subprocess.run(
                ("/bin/ps", "-o", "vsz=", "-p", str(os.getpid())),
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.DEVNULL,
                text=True,
                timeout=2,
            )
            baseline = int(measured.stdout.strip()) * 1024
            if not 1 <= baseline <= (1 << 60) - limit_bytes:
                raise RuntimeError
            absolute_limit = baseline + limit_bytes
        target = (
            absolute_limit
            if inherited_hard == resource.RLIM_INFINITY
            else min(absolute_limit, inherited_hard)
        )
        if target <= 0:
            raise RuntimeError
        # macOS requires lowering the infinite soft limit before its hard limit.
        resource.setrlimit(resource.RLIMIT_AS, (target, inherited_hard))
        resource.setrlimit(resource.RLIMIT_AS, (target, target))
        applied_soft, applied_hard = resource.getrlimit(resource.RLIMIT_AS)
    except (
        AttributeError,
        ImportError,
        OSError,
        OverflowError,
        subprocess.SubprocessError,
        ValueError,
    ):
        raise RuntimeError("office read memory limit is unavailable") from None
    if (
        applied_soft == resource.RLIM_INFINITY
        or applied_hard == resource.RLIM_INFINITY
        or applied_soft != target
        or applied_hard != target
        or target > absolute_limit
    ):
        raise RuntimeError("office read memory limit was not applied")


def _verify_ooxml_archive(content: bytes) -> None:
    try:
        with zipfile.ZipFile(BytesIO(content)) as archive:
            members = archive.infolist()
            if not 1 <= len(members) <= _MAX_OOXML_FILES:
                raise ValueError("OOXML archive file count is invalid")
            total = 0
            seen: set[str] = set()
            for member in members:
                name = member.filename
                relative = PurePosixPath(name)
                if (
                    not name
                    or "\\" in name
                    or relative.is_absolute()
                    or any(part in {"", ".", ".."} for part in relative.parts)
                    or name.casefold() in seen
                    or member.flag_bits & 0x1
                    or member.file_size > _MAX_OOXML_UNCOMPRESSED_BYTES
                ):
                    raise ValueError("OOXML archive member is invalid")
                seen.add(name.casefold())
                total += member.file_size
                if total > _MAX_OOXML_UNCOMPRESSED_BYTES:
                    raise ValueError("OOXML archive expands beyond its limit")
    except zipfile.BadZipFile:
        raise ValueError("OOXML archive is invalid") from None


def _office_read(payload: Mapping[str, Any], runtime: Path) -> Mapping[str, Any]:
    family, content = _office_content(payload)
    collector = _TextCollector()
    warnings = ["visual_layout_not_verified"]
    scan_limited = False

    if family == "document":
        import docx

        _inside(docx, runtime)
        _verify_ooxml_archive(content)
        document = docx.Document(BytesIO(content))
        for paragraph in document.paragraphs:
            if not collector.add(paragraph.text):
                break
        if not collector.truncated:
            for table_index, table in enumerate(document.tables, start=1):
                if not collector.add(f"# Table {table_index}"):
                    break
                for row in table.rows:
                    if not collector.add("\t".join(cell.text for cell in row.cells)):
                        break
                if collector.truncated:
                    break
        structure = {
            "paragraph_count": len(document.paragraphs),
            "table_count": len(document.tables),
        }
    elif family == "spreadsheet":
        import openpyxl

        _inside(openpyxl, runtime)
        _verify_ooxml_archive(content)
        workbook = openpyxl.load_workbook(
            BytesIO(content),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
        observed_formulas = 0
        observed_rows = 0
        try:
            for sheet in workbook.worksheets:
                if not collector.add(f"# Sheet: {sheet.title}"):
                    break
                remaining_rows = 10_000 - observed_rows
                if remaining_rows <= 0:
                    scan_limited = True
                    break
                max_row = min(max(int(sheet.max_row or 1), 1), remaining_rows)
                max_column = min(max(int(sheet.max_column or 1), 1), 256)
                for row in sheet.iter_rows(max_row=max_row, max_col=max_column):
                    observed_rows += 1
                    values = []
                    for cell in row:
                        if getattr(cell, "data_type", None) == "f":
                            observed_formulas += 1
                        value = cell.value
                        values.append("" if value is None else str(value))
                    if not collector.add("\t".join(values).rstrip("\t")):
                        break
                if int(sheet.max_row or 0) > max_row or int(sheet.max_column or 0) > 256:
                    scan_limited = True
                if collector.truncated or scan_limited:
                    break
            structure = {
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": list(workbook.sheetnames[:100]),
                "observed_row_count": observed_rows,
                "observed_formula_count": observed_formulas,
            }
        finally:
            workbook.close()
        if observed_formulas:
            warnings.append("formulas_not_calculated")
    elif family == "presentation":
        import pptx

        _inside(pptx, runtime)
        _verify_ooxml_archive(content)
        presentation = pptx.Presentation(BytesIO(content))
        for slide_index, slide in enumerate(presentation.slides, start=1):
            if not collector.add(f"# Slide {slide_index}"):
                break
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and not collector.add(shape.text):
                    break
            if collector.truncated:
                break
        structure = {"slide_count": len(presentation.slides)}
    else:
        import pypdf

        _inside(pypdf, runtime)
        reader = pypdf.PdfReader(BytesIO(content), strict=True)
        if reader.is_encrypted:
            raise ValueError("encrypted PDF cannot be inspected")
        pages_without_text = 0
        for page_index, page in enumerate(reader.pages, start=1):
            if page_index > 500:
                scan_limited = True
                break
            page_text = page.extract_text() or ""
            if not page_text.strip():
                pages_without_text += 1
            if not collector.add(f"# Page {page_index}\n{page_text}"):
                break
        structure = {
            "page_count": len(reader.pages),
            "pages_without_extractable_text": pages_without_text,
        }
        if pages_without_text:
            warnings.append("some_pages_have_no_extractable_text")

    if scan_limited:
        warnings.append("structure_scan_limited")
    return {
        "provider": "python-office-formats-v1",
        "family": family,
        "text": collector.value(),
        "structure": structure,
        "warnings": warnings,
        "truncated": collector.truncated or scan_limited,
    }


def _office_create(payload: Mapping[str, Any], runtime: Path) -> Mapping[str, Any]:
    family = str(payload.get("family") or "")
    title = _text(payload.get("title") or "e-Mate 办公产物", maximum=512)
    body = BytesIO()
    validation: dict[str, Any]

    if family == "document":
        import docx

        _inside(docx, runtime)
        document = docx.Document()
        document.add_heading(title, level=0)
        sections = payload.get("sections") or []
        for section in sections:
            if not isinstance(section, Mapping):
                raise ValueError("document section is invalid")
            heading = str(section.get("heading") or "").strip()
            if heading:
                document.add_heading(_text(heading, maximum=512), level=1)
            paragraphs = section.get("paragraphs") or []
            if not isinstance(paragraphs, list):
                raise ValueError("document paragraphs are invalid")
            for paragraph in paragraphs:
                document.add_paragraph(_text(paragraph))
        document.save(body)
        reopened = docx.Document(BytesIO(body.getvalue()))
        validation = {"paragraph_count": len(reopened.paragraphs)}
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        extension = ".docx"
    elif family == "spreadsheet":
        import openpyxl

        _inside(openpyxl, runtime)
        workbook = openpyxl.Workbook()
        workbook.remove(workbook.active)
        sheets = payload.get("sheets") or []
        if not isinstance(sheets, list) or not sheets:
            raise ValueError("spreadsheet sheets are invalid")
        for sheet_payload in sheets:
            if not isinstance(sheet_payload, Mapping):
                raise ValueError("spreadsheet sheet is invalid")
            sheet = workbook.create_sheet(_text(sheet_payload.get("name"), maximum=64))
            rows = sheet_payload.get("rows") or []
            if not isinstance(rows, list):
                raise ValueError("spreadsheet rows are invalid")
            for row in rows:
                if not isinstance(row, list):
                    raise ValueError("spreadsheet row is invalid")
                sheet.append(row)
        workbook.save(body)
        reopened = openpyxl.load_workbook(BytesIO(body.getvalue()), read_only=True)
        validation = {"sheet_count": len(reopened.sheetnames)}
        reopened.close()
        mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        extension = ".xlsx"
    elif family == "presentation":
        import pptx

        _inside(pptx, runtime)
        presentation = pptx.Presentation()
        slides = payload.get("slides") or []
        if not isinstance(slides, list) or not slides:
            raise ValueError("presentation slides are invalid")
        for index, slide_payload in enumerate(slides):
            if not isinstance(slide_payload, Mapping):
                raise ValueError("presentation slide is invalid")
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[0 if index == 0 else 1]
            )
            slide.shapes.title.text = _text(slide_payload.get("title"), maximum=512)
            bullets = slide_payload.get("bullets") or []
            if not isinstance(bullets, list):
                raise ValueError("presentation bullets are invalid")
            if index == 0:
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = "\n".join(_text(item) for item in bullets)
            else:
                frame = slide.placeholders[1].text_frame
                frame.clear()
                for bullet_index, bullet in enumerate(bullets):
                    paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                    paragraph.text = _text(bullet)
        presentation.save(body)
        reopened = pptx.Presentation(BytesIO(body.getvalue()))
        validation = {"slide_count": len(reopened.slides)}
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        extension = ".pptx"
    elif family == "pdf":
        import pypdf
        import reportlab
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.pdfgen import canvas

        for module in (pypdf, reportlab):
            _inside(module, runtime)
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        writer = canvas.Canvas(body, pagesize=A4)
        width, height = A4
        y = height - 64
        writer.setFont("STSong-Light", 20)
        writer.drawString(56, y, title)
        y -= 36
        writer.setFont("STSong-Light", 11)
        for section in payload.get("sections") or []:
            if not isinstance(section, Mapping):
                raise ValueError("PDF section is invalid")
            heading = str(section.get("heading") or "").strip()
            lines = ([heading] if heading else []) + list(section.get("paragraphs") or [])
            for value in lines:
                for line in textwrap.wrap(_text(value), width=55, replace_whitespace=False):
                    if y < 56:
                        writer.showPage()
                        writer.setFont("STSong-Light", 11)
                        y = height - 56
                    writer.drawString(56, y, line)
                    y -= 18
                y -= 6
        writer.save()
        reopened = pypdf.PdfReader(BytesIO(body.getvalue()))
        validation = {"page_count": len(reopened.pages)}
        mime_type = "application/pdf"
        extension = ".pdf"
    else:
        raise ValueError("office family is unsupported")

    content = body.getvalue()
    # Base64 and the response envelope must remain under the process adapter's
    # 8 MiB stdout boundary.
    if not 1 <= len(content) <= 5 * 1024 * 1024:
        raise ValueError("office output size is invalid")
    return {
        "provider": "python-office-formats-v1",
        "family": family,
        "mime_type": mime_type,
        "extension": extension,
        "size_bytes": len(content),
        "content_base64": base64.b64encode(content).decode("ascii"),
        "validation": validation,
    }


def _office_edit(payload: Mapping[str, Any], runtime: Path) -> Mapping[str, Any]:
    """Replace one validated Office artifact with complete structured content."""

    family, _content = _office_content(payload)
    # Opening the source before replacement proves this is an edit of the
    # supplied artifact, not a create call mislabeled by the transport.
    _office_read(payload, runtime)
    result = dict(_office_create(payload, runtime))
    result["validation"] = {
        **dict(result["validation"]),
        "source_opened": True,
    }
    return result


def main() -> int:
    if len(sys.argv) not in {2, 4}:
        return 2
    memory_limit: int | None = None
    if len(sys.argv) == 4:
        if sys.argv[2] != "--office-read-memory-limit":
            return 2
        try:
            memory_limit = int(sys.argv[3])
        except ValueError:
            return 2
    root = Path(sys.argv[1]).resolve(strict=True)
    runtime = (root / "runtime" / "python").resolve(strict=True)
    runtime.relative_to(root)
    sys.path.insert(0, str(runtime))
    try:
        request = json.loads(
            sys.stdin.buffer.read(12 * 1024 * 1024 + 1).decode("utf-8")
        )
        pack_id = request.get("pack_id")
        operation = request.get("operation")
        payload = request.get("payload")
        if request.get("schema_version") != 1 or not isinstance(payload, Mapping):
            return 3
        if pack_id == "ocr" and operation == "extract":
            result = _ocr(payload, runtime)
        elif pack_id == "office" and operation == "probe":
            result = _office_probe(runtime)
        elif pack_id == "office" and operation == "create":
            result = _office_create(payload, runtime)
        elif pack_id == "office" and operation == "read":
            if memory_limit is None:
                return 4
            _apply_office_read_memory_limit(memory_limit)
            result = _office_read(payload, runtime)
        elif pack_id == "office" and operation == "edit":
            if memory_limit is None:
                return 4
            _apply_office_read_memory_limit(memory_limit)
            result = _office_edit(payload, runtime)
        else:
            return 4
        sys.stdout.write(
            json.dumps(
                {
                    "schema_version": 1,
                    "pack_id": pack_id,
                    "status": "success",
                    "result": result,
                },
                sort_keys=True,
                separators=(",", ":"),
            )
        )
        return 0
    except Exception:
        return 5


if __name__ == "__main__":
    raise SystemExit(main())
