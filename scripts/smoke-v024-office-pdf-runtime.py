#!/usr/bin/env python3
"""Smoke-test the shared v0.2.4 Office/PDF runtime primitives."""

from __future__ import annotations

import argparse
import json
import pathlib
import sys
import tempfile
from datetime import datetime, timezone
from typing import Any


ROOT = pathlib.Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from common.office_pdf_runtime import (  # noqa: E402
    ARTIFACT_KINDS,
    build_quality_evidence,
    default_quality_gates,
    inspect_office_pdf_artifact,
    probe_office_pdf_runtime,
    render_pdf_pages,
)


def add_check(checks: list[dict[str, Any]], label: str, ok: bool, evidence: Any) -> None:
    checks.append({
        "label": label,
        "status": "PASS" if ok else "FAIL",
        "evidence": evidence,
    })


def add_warning(checks: list[dict[str, Any]], label: str, evidence: Any) -> None:
    checks.append({"label": label, "status": "WARN", "evidence": evidence})


def create_docx(path: pathlib.Path) -> bool:
    try:
        import docx  # type: ignore
    except Exception:
        return False
    document = docx.Document()
    document.add_heading("Private smoke title", level=1)
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Private cell"
    document.save(path)
    return True


def create_xlsx(path: pathlib.Path) -> bool:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return False
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A1"] = "Private value"
    sheet["B1"] = "=SUM(1,2)"
    workbook.save(path)
    return True


def create_pptx(path: pathlib.Path) -> bool:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return False
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Private deck title"
    presentation.save(path)
    return True


def create_pdf(path: pathlib.Path) -> bool:
    try:
        from reportlab.pdfgen import canvas  # type: ignore
    except Exception:
        return False
    c = canvas.Canvas(str(path))
    c.drawString(72, 720, "Private PDF text")
    c.showPage()
    c.save()
    return True


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", default="")
    parser.add_argument("--require-pdf-render", action="store_true")
    args = parser.parse_args(argv)

    checks: list[dict[str, Any]] = []
    probe = probe_office_pdf_runtime()
    add_check(checks, "office-pdf runtime probe schema", probe.get("schemaVersion") == 1 and probe.get("redacted") is True, {"status": probe.get("status")})

    for kind, spec in ARTIFACT_KINDS.items():
        row = probe["artifactKinds"].get(kind) or {}
        add_check(
            checks,
            f"{kind} parser/writer readiness",
            row.get("parseStatus") == "ready" and row.get("writeStatus") in {"ready", "not-applicable"},
            {
                "parseStatus": row.get("parseStatus"),
                "writeStatus": row.get("writeStatus"),
                "missingParserModules": row.get("missingParserModules"),
                "missingWriterModules": row.get("missingWriterModules"),
            },
        )
        add_check(
            checks,
            f"{kind} quality gate contract",
            row.get("qualityGates") == spec.get("qualityGates"),
            {"qualityGates": row.get("qualityGates")},
        )

    samples: list[dict[str, Any]] = []
    with tempfile.TemporaryDirectory(prefix="ecorex-v024-office-pdf-") as tmp:
        root = pathlib.Path(tmp)
        sample_specs = [
            ("document", root / "sample.docx", create_docx),
            ("spreadsheet", root / "sample.xlsx", create_xlsx),
            ("presentation", root / "sample.pptx", create_pptx),
            ("pdf", root / "sample.pdf", create_pdf),
        ]
        for kind, path, factory in sample_specs:
            created = factory(path)
            add_check(checks, f"{kind} sample creation", created, {"extension": path.suffix})
            if not created:
                continue
            inspected = inspect_office_pdf_artifact(path)
            samples.append({
                "kind": kind,
                "extension": inspected.get("extension"),
                "summary": inspected.get("summary"),
                "checks": inspected.get("checks"),
            })
            add_check(checks, f"{kind} content-free inspection", inspected.get("kind") == kind and bool(inspected.get("summary")), {"summaryKeys": sorted((inspected.get("summary") or {}).keys())})

        pdf_path = root / "sample.pdf"
        if pdf_path.exists() and probe["artifactKinds"]["pdf"].get("renderStatus") == "ready":
            rendered = render_pdf_pages(pdf_path, root / "pdf-pages", max_pages=1, dpi=96)
            add_check(
                checks,
                "pdf page render primitive",
                rendered.get("renderedCount") == 1 and (rendered.get("artifacts") or [{}])[0].get("width", 0) > 0,
                {
                    "renderedCount": rendered.get("renderedCount"),
                    "artifactCount": len(rendered.get("artifacts") or []),
                    "firstArtifact": (rendered.get("artifacts") or [{}])[0],
                },
            )
        elif args.require_pdf_render:
            add_check(checks, "pdf page render primitive", False, {"renderStatus": probe["artifactKinds"]["pdf"].get("renderStatus")})
        else:
            add_warning(checks, "pdf page render primitive", {"renderStatus": probe["artifactKinds"]["pdf"].get("renderStatus")})

    for kind in ARTIFACT_KINDS:
        gates = default_quality_gates(kind)
        evidence = build_quality_evidence(kind, [{"id": gate, "status": "pass"} for gate in gates])
        add_check(checks, f"{kind} shared QA evidence schema", evidence.get("status") == "pass" and evidence.get("redacted") is True, {"gateCount": len(gates)})

    failures = [item for item in checks if item["status"] == "FAIL"]
    payload = {
        "status": "PASS" if not failures else "FAIL",
        "generatedAt": datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z"),
        "scope": "R24-04 shared Office/PDF runtime primitives",
        "probeSummary": {
            "status": probe.get("status"),
            "pdfRenderStatus": probe["artifactKinds"]["pdf"].get("renderStatus"),
            "artifactToolAvailable": bool(probe.get("artifactToolAvailable")),
            "kindStatuses": {
                kind: {
                    "parseStatus": row.get("parseStatus"),
                    "writeStatus": row.get("writeStatus"),
                    "renderStatus": row.get("renderStatus"),
                }
                for kind, row in probe.get("artifactKinds", {}).items()
            },
        },
        "sampleInspections": samples,
        "checks": checks,
        "failed": failures,
        "redacted": True,
    }
    text = json.dumps(payload, ensure_ascii=False, indent=2) + "\n"
    if args.output:
        output = pathlib.Path(args.output)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(text, encoding="utf-8")
    print(text, end="")
    return 0 if not failures else 1


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
