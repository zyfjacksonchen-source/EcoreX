from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path
from typing import Any, Dict

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from common.office_pdf_runtime import (
    OfficePdfRuntimeError,
    build_presentation_quality_evidence,
    render_presentation_preview,
)


def _write_json(path: Path, payload: Dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def _render_metadata(slide_count: int) -> list[Dict[str, Any]]:
    return [
        {"slide": slide, "artifactRef": f"hmac:presentation-render-{slide:03d}", "width": 1280, "height": 720}
        for slide in range(1, slide_count + 1)
    ]


def _create_clean_deck(path: Path) -> None:
    import pptx
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Private R24 presentation story"
    title_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(54)
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Sensitive launch plan"
    subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(28)

    chart_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_slide.shapes.title.text = "Quality trend"
    chart_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(38)
    body = chart_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.0), Inches(1.2))
    body.text_frame.text = "Private metric summary"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Hidden Series", (1, 2, 3))
    chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5.1),
        Inches(1.4),
        Inches(4.2),
        Inches(3.9),
        chart_data,
    )
    presentation.save(path)


def _create_bad_deck(path: Path) -> None:
    import pptx
    from pptx.util import Inches, Pt

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Tiny private title"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    first = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(3.0), Inches(1.2))
    first.text_frame.text = "Private overlap one"
    first.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    second = slide.shapes.add_textbox(Inches(1.3), Inches(1.7), Inches(3.0), Inches(1.2))
    second.text_frame.text = "Private overlap two"
    second.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
    clipped = slide.shapes.add_textbox(Inches(12.8), Inches(7.2), Inches(1.0), Inches(1.0))
    clipped.text = "Private clipped"
    presentation.save(path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Smoke R24-05 presentation quality evidence.")
    parser.add_argument("--output", default="docs/v0.2.4/artifacts/presentation-quality-smoke.json")
    args = parser.parse_args()

    with tempfile.TemporaryDirectory(prefix="ecorex-ppt-quality-") as tmp:
        root = Path(tmp)
        clean_path = root / "clean-private-deck.pptx"
        bad_path = root / "bad-private-deck.pptx"
        _create_clean_deck(clean_path)
        _create_bad_deck(bad_path)

        render_backend = "provided-redacted-metadata"
        actual_render = None
        actual_render_available = False
        try:
            actual_render = render_presentation_preview(clean_path, root / "rendered", max_slides=2)
            render_backend = str(actual_render.get("renderBackend") or "libreoffice-pdf")
            render_items = actual_render.get("artifacts") or []
            actual_render_available = True
        except OfficePdfRuntimeError as exc:
            render_items = _render_metadata(2)
            actual_render = {"status": "unavailable", "errorType": type(exc).__name__}

        clean = build_presentation_quality_evidence(
            clean_path,
            authoring_route="artifact-tool",
            renders=render_items,
            visual_inspection_passed=True,
        )
        bad = build_presentation_quality_evidence(
            bad_path,
            authoring_route="",
            renders=[],
            visual_inspection_passed=False,
        )
        required_bad_failures = {"layout-bounds", "font-size-check", "render-preview", "overlap-check"}
        bad_failed_checks = {item["id"] for item in bad.get("checks", []) if item.get("status") == "fail"}
        clean_failed_checks = {item["id"] for item in clean.get("checks", []) if item.get("status") == "fail"}
        synthetic_render_rejected = (
            not actual_render_available
            and "render-preview" in clean_failed_checks
            and not clean.get("renderedArtifacts")
        )
        clean_render_contract_ok = (
            clean.get("status") == "pass"
            if actual_render_available
            else clean.get("status") == "fail" and synthetic_render_rejected
        )
        serialized = json.dumps({"clean": clean, "bad": bad}, ensure_ascii=False)
        leaks = [
            item
            for item in (
                "Private R24 presentation story",
                "Sensitive launch plan",
                "Private metric summary",
                "Private overlap",
                str(root),
                clean_path.name,
                bad_path.name,
            )
            if item and item in serialized
        ]
        payload = {
            "status": (
                "PASS"
                if (
                    clean_render_contract_ok
                    and bad.get("status") == "fail"
                    and required_bad_failures <= bad_failed_checks
                    and not leaks
                )
                else "FAIL"
            ),
            "renderBackend": render_backend,
            "actualRender": actual_render,
            "actualRenderAvailable": actual_render_available,
            "syntheticRenderRejected": synthetic_render_rejected,
            "cleanStatus": clean.get("status"),
            "cleanFailedChecks": sorted(clean_failed_checks),
            "badStatus": bad.get("status"),
            "badFailedChecks": sorted(bad_failed_checks),
            "requiredBadFailures": sorted(required_bad_failures),
            "cleanSummary": clean.get("presentationAnalysis", {}).get("summary", {}),
            "badSummary": bad.get("presentationAnalysis", {}).get("summary", {}),
            "leakCount": len(leaks),
            "redacted": True,
        }
        _write_json(Path(args.output), payload)
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
